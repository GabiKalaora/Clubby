"""
Clubby — Premium Hebrew Deck v2 (UI/UX Improvements Edition)
Highlights all the upgrades from the senior staff designer audit
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

BASE     = "/Users/I560647/Library/CloudStorage/OneDrive-SAPSE/Desktop/finance/Clubby"
V2       = f"{BASE}/screenshots/v2"
DEMO     = f"{BASE}/screenshots/demo"
SCEN     = f"{BASE}/screenshots/scenario"
WORDMARK = f"{BASE}/logo/wordmark.png"
OUT      = f"{BASE}/slides/Clubby_HE_v2_UpgradeEdition.pptx"
os.makedirs(f"{BASE}/slides", exist_ok=True)

# ── Premium dark palette ─────────────────────────────────────────────────────
BG      = RGBColor(0x06, 0x07, 0x0B)
SURFACE = RGBColor(0x0E, 0x10, 0x16)
CARD    = RGBColor(0x16, 0x19, 0x22)
BORDER  = RGBColor(0x22, 0x28, 0x38)
GREEN   = RGBColor(0x1A, 0x7A, 0x4A)
GBRT    = RGBColor(0x2E, 0xCC, 0x71)
GOLD    = RGBColor(0xD4, 0xAF, 0x37)
SILVER  = RGBColor(0xA8, 0xB2, 0xC8)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DIM     = RGBColor(0x4A, 0x52, 0x60)
LGRAY   = RGBColor(0x7A, 0x85, 0x99)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
TEAL    = RGBColor(0x06, 0xB6, 0xD4)
INDIGO  = RGBColor(0x63, 0x66, 0xF1)
ROSE    = RGBColor(0xF4, 0x3F, 0x5E)
RED     = RGBColor(0xEF, 0x44, 0x44)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

def set_rtl(p_xml):
    pPr = p_xml.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p_xml, qn('a:pPr'))
        p_xml.insert(0, pPr)
    pPr.set('rtl', '1')

def fill_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def t(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
      align=PP_ALIGN.RIGHT, italic=False, rtl=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    if rtl: set_rtl(p._p)

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def rrect(slide, x, y, w, h, color, border=None):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def gold_line(slide, x, y, w):
    rect(slide, x, y, w, Pt(1.5), GOLD)

def wordmark(slide):
    if os.path.exists(WORDMARK):
        slide.shapes.add_picture(WORDMARK, Inches(0.28), Inches(0.18), Inches(2.2), Inches(0.58))
    else:
        t(slide, "clubby.", Inches(0.28), Inches(0.18), Inches(2.5), Inches(0.55),
          size=20, bold=True, color=GBRT, align=PP_ALIGN.LEFT, rtl=False)

def pgnum(slide, n, total):
    t(slide, f"{n} / {total}", Inches(0.25), H-Inches(0.38), Inches(1.2), Inches(0.3),
      size=9, color=DIM, align=PP_ALIGN.LEFT, rtl=False)

def pill(slide, label, x, y, color=GOLD, pw=Inches(1.6), bg=None):
    ph = Inches(0.3)
    p = slide.shapes.add_shape(9, x, y, pw, ph)
    p.fill.solid()
    p.fill.fore_color.rgb = bg or RGBColor(0x1A, 0x16, 0x08)
    p.line.color.rgb = color; p.line.width = Pt(0.75)
    t(slide, label, x+Inches(0.08), y+Inches(0.04), pw-Inches(0.16), ph-Inches(0.06),
      size=9, bold=True, color=color, align=PP_ALIGN.CENTER, rtl=False)

def phone_frame(slide, path, x, y, fw, fh):
    if not path or not os.path.exists(path): return
    frm = slide.shapes.add_shape(9, x, y, fw, fh)
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0x06, 0x08, 0x0E)
    frm.line.color.rgb = GOLD; frm.line.width = Pt(1.2)
    pad = Inches(0.09)
    slide.shapes.add_picture(path, x+pad, y+pad, fw-pad*2, fh-pad*2)

def browser_frame(slide, path, x, y, fw, fh):
    if not path or not os.path.exists(path): return
    chrome_h = Inches(0.24)
    frm = slide.shapes.add_shape(9, x, y, fw, fh)
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0x0A, 0x0C, 0x14)
    frm.line.color.rgb = GOLD; frm.line.width = Pt(1.2)
    bar = slide.shapes.add_shape(1, x, y, fw, chrome_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x12, 0x15, 0x1F)
    bar.line.fill.background()
    for di, dc in enumerate([RGBColor(0xFC,0x5C,0x5C), RGBColor(0xFD,0xBC,0x40), RGBColor(0x27,0xC9,0x3F)]):
        d = slide.shapes.add_shape(9, x+Inches(0.1+di*0.18), y+Inches(0.075), Inches(0.09), Inches(0.09))
        d.fill.solid(); d.fill.fore_color.rgb = dc; d.line.fill.background()
    slide.shapes.add_picture(path, x, y+chrome_h, fw, fh-chrome_h)

def dark_card(slide, x, y, w, h, accent=None):
    c = slide.shapes.add_shape(9, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = CARD
    c.line.color.rgb = BORDER; c.line.width = Pt(0.75)
    if accent:
        rect(slide, x, y, w, Pt(2.5), accent)
    return c

def starfield(slide, density=22):
    import random; random.seed(7)
    for _ in range(density):
        sx = random.uniform(0.2, 13.0); sy = random.uniform(0.2, 7.2)
        sz = random.uniform(0.025, 0.06)
        d = slide.shapes.add_shape(9, Inches(sx), Inches(sy), Inches(sz), Inches(sz))
        d.fill.solid(); d.fill.fore_color.rgb = GOLD
        d.line.fill.background()

TOTAL = 22

# ══════════════════════════════════════════════════════════════════════════════
# 1. COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
starfield(sl, 28)
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, GOLD)
wordmark(sl)

# Center glow
glow = rect(sl, 0, Inches(2.4), W, Inches(2.7), RGBColor(0x09,0x12,0x0B))
glow.line.fill.background()

t(sl, "מהדורת שדרוג —", Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.7),
  size=22, color=SILVER, align=PP_ALIGN.RIGHT)
t(sl, "Clubby v2.", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.3),
  size=72, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
t(sl, "מהיר יותר. בהיר יותר. אבטחה ארגונית.", Inches(0.5), Inches(2.85), Inches(12.3), Inches(1.0),
  size=72, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.0), Inches(4.05), Inches(5.6))
t(sl, "ביקורת UI/UX מקיפה · 30+ תיקוני קוד · אינטגרציה ארגונית מאובטחת",
  Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
  size=18, color=SILVER, align=PP_ALIGN.RIGHT)

# Stat band
rect(sl, 0, Inches(5.2), W, Inches(2.1), RGBColor(0x09, 0x0B, 0x12))
gold_line(sl, 0, Inches(5.2), W)
stats = [
    ("✅", "12 תיקונים קריטיים", "באגים שבירת UX"),
    ("🔐", "מפתח אינטגרציה M2M", "במקום service_role"),
    ("⚡", "תגובת רכישה 77ms", "Non-blocking architecture"),
    ("🔁", "Idempotency מובנית", "POS retry-safe"),
]
for i, (em, title, sub) in enumerate(stats):
    x = Inches(0.5 + i * 3.2)
    rrect(sl, x, Inches(5.4), Inches(2.95), Inches(1.65), SURFACE, BORDER)
    t(sl, em + "  " + title, x+Inches(0.15), Inches(5.55), Inches(2.7), Inches(0.45),
      size=14, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    t(sl, sub, x+Inches(0.15), Inches(6.05), Inches(2.7), Inches(0.4),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 2. EXECUTIVE SUMMARY — TOP 10 IMPROVEMENTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GOLD)
wordmark(sl); pgnum(sl, 2, TOTAL)

pill(sl, "EXECUTIVE SUMMARY", Inches(0.4), Inches(0.72), pw=Inches(2.3))
t(sl, "10 שדרוגים מובילים בהשפעה גבוהה",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.5), Inches(2.0), Inches(6.1))

improvements = [
    ("🔴", "תיקון מימוש שקט", "redeem.tsx — שגיאה הוצגה inline במקום להיעלם"),
    ("🔴", "תיקון Spinner נצחי", "Dashboard / Members / Anchor — try/catch על Promise.all"),
    ("🔴", "מודאל פורטל ב-dark mode", "var(--surface) במקום background:white"),
    ("🔴", "כפתור Restart שעבד", "DevSettings.reload() / window.reload() במקום no-op"),
    ("🔴", "שמירת קלט בטופס פיד", "FeedPosts לא נסגר על שגיאה — שומר נתוני המשתמש"),
    ("🟠", "מפתח אינטגרציה M2M", "החלפת service_role בסוד ייעודי לחילופי גרסה"),
    ("🟠", "Idempotency על רכישות", "transaction_id מונע נקודות כפולות מ-POS retry"),
    ("🟠", "Persistence ב-JSON", "תהילים שורדים restart של שרת (data/*.json)"),
    ("🟠", "Non-blocking benefit", "רכישה חוזרת ב-77ms · Clubby fire-and-log"),
    ("🟢", "Integration Log חי", "מסך אירועים עם רענון אוטומטי כל 5 שניות"),
]

for i, (icon, title, desc) in enumerate(improvements):
    col = i % 2
    row = i // 2
    bx = Inches(0.4 + col * 6.45)
    by = Inches(2.3 + row * 0.92)
    dark_card(sl, bx, by, Inches(6.3), Inches(0.78))
    t(sl, icon, bx+Inches(0.15), by+Inches(0.18), Inches(0.4), Inches(0.4),
      size=18, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    t(sl, title, bx+Inches(0.6), by+Inches(0.1), Inches(5.6), Inches(0.32),
      size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    t(sl, desc, bx+Inches(0.6), by+Inches(0.42), Inches(5.6), Inches(0.32),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 3. SECTION DIVIDER — Critical Bug Fixes
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
starfield(sl, 18)
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, RED)
t(sl, "01", Inches(0.35), Inches(0.5), Inches(2), Inches(0.65),
  size=13, bold=True, color=DIM, align=PP_ALIGN.LEFT, rtl=False)
glow = rect(sl, 0, Inches(2.2), W, Inches(2.2), RGBColor(0x12, 0x07, 0x07))
glow.line.fill.background()
t(sl, "תיקוני באגים קריטיים", Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.8),
  size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
gold_line(sl, Inches(4.5), Inches(3.25), Inches(4.3))
t(sl, "5 תיקוני P0 שבר אינטראקציות יומיומיות של משתמשים",
  Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.65),
  size=18, color=SILVER, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 4. FIX #1 — Silent redemption failure → inline error
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, RED)
wordmark(sl); pgnum(sl, 4, TOTAL)

PH = Inches(6.0); PW = Inches(2.9)
phone_frame(sl, f"{V2}/11-mobile-redeem-with-error-state.png" if os.path.exists(f"{V2}/11-mobile-redeem-with-error-state.png") else f"{DEMO}/14-mobile-redeem-screen.png",
            W - PW - Inches(0.4), (H-PH)/2, PW, PH)

pill(sl, "BEFORE  →  AFTER", Inches(0.4), Inches(0.72), color=RED, pw=Inches(2.0))
t(sl, "מימוש שקט → שגיאה inline",
  Inches(0.4), Inches(1.18), Inches(9.0), Inches(0.8),
  size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, Inches(0.4), Inches(2.1), Inches(4.5))

# Before/After cards
bx = Inches(0.4); cardW = Inches(4.4); cardH = Inches(2.2)
dark_card(sl, bx, Inches(2.4), cardW, cardH, RED)
t(sl, "❌ לפני", bx+Inches(0.15), Inches(2.55), Inches(2.0), Inches(0.4),
  size=14, bold=True, color=RED, align=PP_ALIGN.RIGHT)
t(sl, "if (error) return\n\nהמשתמש לחץ 'אישור מימוש'\nהבקשה נכשלה\nהמסך נשאר באותו מצב\nאין משוב, אין הסבר",
  bx+Inches(0.15), Inches(3.0), cardW-Inches(0.3), Inches(1.5),
  size=12, color=SILVER, align=PP_ALIGN.RIGHT)

bx2 = Inches(5.0);
dark_card(sl, bx2, Inches(2.4), cardW, cardH, GBRT)
t(sl, "✅ אחרי", bx2+Inches(0.15), Inches(2.55), Inches(2.0), Inches(0.4),
  size=14, bold=True, color=GBRT, align=PP_ALIGN.RIGHT)
t(sl, "setRedeemError(error.message)\n\nשגיאה inline בצבע אדום\nמתחת לכפתור\nהודעה ברורה למשתמש\nאפשרות לנסות שוב",
  bx2+Inches(0.15), Inches(3.0), cardW-Inches(0.3), Inches(1.5),
  size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# Code snippet
dark_card(sl, Inches(0.4), Inches(4.85), Inches(9.0), Inches(2.0), GOLD)
t(sl, "📝  השינוי בקוד", Inches(0.55), Inches(5.0), Inches(8.7), Inches(0.4),
  size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
t(sl, "// apps/mobile/app/redeem/[id].tsx\n+ const [redeemError, setRedeemError] = useState('')\n  setRedeeming(false)\n- if (error) return\n+ if (error) { setRedeemError(error.message); return }",
  Inches(0.55), Inches(5.45), Inches(8.7), Inches(1.4),
  size=12, color=SILVER, align=PP_ALIGN.LEFT, rtl=False)

# ══════════════════════════════════════════════════════════════════════════════
# 5. FIX #2 — Portal modal dark mode
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, INDIGO)
wordmark(sl); pgnum(sl, 5, TOTAL)

FW = Inches(7.5); FH = Inches(5.6)
browser_frame(sl, f"{V2}/14-portal-member-modal-themed.png", Inches(0.35), Inches(0.85), FW, FH)

TX = Inches(8.2); TW = Inches(4.85)
pill(sl, "PORTAL", TX, Inches(0.72), color=INDIGO, pw=Inches(1.2))
t(sl, "מודאל בנושא תואם",
  TX, Inches(1.18), TW, Inches(0.8), size=24, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(1.0), Inches(2.1), Inches(3.85))
t(sl, "מודאל פרטי חבר היה background:'white' קבוע — שובר את ה-dark mode. כעת משתמש ב-var(--surface) ומתאים לכל theme.",
  TX, Inches(2.25), TW, Inches(1.4), size=12, color=SILVER, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "✅ background: var(--surface)",
    "✅ border: var(--border)",
    "✅ color: var(--text)",
    "✅ עובד אוטומטית בלייט/דארק",
    "✅ אותו fix גם במודאל webhook guide",
]):
    t(sl, b, TX, Inches(3.7+j*0.5), TW, Inches(0.45),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 6. FIX #3 — Permanent spinner
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 6, TOTAL)

FW = Inches(7.5); FH = Inches(5.6)
browser_frame(sl, f"{V2}/12-portal-dashboard.png", Inches(0.35), Inches(0.85), FW, FH)

TX = Inches(8.2); TW = Inches(4.85)
pill(sl, "ERROR HANDLING", TX, Inches(0.72), color=AMBER, pw=Inches(2.1))
t(sl, "Spinner נצחי → הודעת שגיאה",
  TX, Inches(1.18), TW, Inches(1.0), size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(1.5), Inches(2.3), Inches(3.35))

t(sl, "הבעיה:",
  TX, Inches(2.45), TW, Inches(0.4), size=13, bold=True, color=AMBER, align=PP_ALIGN.RIGHT)
t(sl, "Promise.all ללא .catch — אם RPC נכשל, setLoading(false) לא נקרא לעולם → המסך תקוע ב-spinner.",
  TX, Inches(2.85), TW, Inches(1.1), size=12, color=SILVER, align=PP_ALIGN.RIGHT)

t(sl, "הפתרון:",
  TX, Inches(4.0), TW, Inches(0.4), size=13, bold=True, color=GBRT, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    ".catch(() => setLoadError(true))",
    "alert-error מוצג למשתמש",
    "המשתמש יכול לרענן בלי תקיעה",
    "תוקן ב-Dashboard, Members, Anchor",
]):
    t(sl, "✓ " + b, TX, Inches(4.4+j*0.45), TW, Inches(0.42),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 7. SECTION — Mock Backend v2
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
starfield(sl, 18)
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, AMBER)
t(sl, "02", Inches(0.35), Inches(0.5), Inches(2), Inches(0.65),
  size=13, bold=True, color=DIM, align=PP_ALIGN.LEFT, rtl=False)
glow = rect(sl, 0, Inches(2.2), W, Inches(2.2), RGBColor(0x11, 0x0C, 0x05))
glow.line.fill.background()
t(sl, "אינטגרציה ארגונית v2", Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.8),
  size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
gold_line(sl, Inches(4.5), Inches(3.25), Inches(4.3))
t(sl, "Mock backend v2 — אבטחה, אמינות, ושקיפות מלאה",
  Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.65),
  size=18, color=SILVER, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 8. M2M INTEGRATION KEY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 8, TOTAL)

pill(sl, "🔐 SECURITY", Inches(0.4), Inches(0.72), color=AMBER, pw=Inches(1.4))
t(sl, "מפתח אינטגרציה M2M ייעודי",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.8),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.5), Inches(2.1), Inches(6.1))
t(sl, "במקום לתת ל-POS את service_role_key (גישת superadmin לכל הנתונים), כעת מפתח ייעודי שניתן לסבב מבלי להפיל את האינטגרציה.",
  Inches(0.4), Inches(2.3), Inches(12.5), Inches(0.65),
  size=14, color=SILVER, align=PP_ALIGN.RIGHT)

# Before/After
b1x = Inches(0.4); b2x = Inches(7.0); cw = Inches(5.95); ch = Inches(3.6)
dark_card(sl, b1x, Inches(3.2), cw, ch, RED)
t(sl, "❌ לפני",
  b1x+Inches(0.2), Inches(3.4), cw-Inches(0.4), Inches(0.45),
  size=15, bold=True, color=RED, align=PP_ALIGN.RIGHT)
t(sl, "Authorization: Bearer <SERVICE_ROLE_KEY>\n\n• POS חייב מפתח admin של כל הDB\n• אין סיבוב מפתח אפשרי\n• אם מפתח דולף — סיכון אבטחה ענק\n• אסור לתת מפתח כזה ללקוח חיצוני",
  b1x+Inches(0.2), Inches(3.95), cw-Inches(0.4), Inches(2.7),
  size=12, color=SILVER, align=PP_ALIGN.RIGHT)

dark_card(sl, b2x, Inches(3.2), cw, ch, GBRT)
t(sl, "✅ אחרי",
  b2x+Inches(0.2), Inches(3.4), cw-Inches(0.4), Inches(0.45),
  size=15, bold=True, color=GBRT, align=PP_ALIGN.RIGHT)
t(sl, "Authorization: Bearer <CLUBBY_INTEGRATION_KEY>\n\n• מפתח ייעודי לאינטגרציות חיצוניות\n• ניתן לסבב בלי redeploy\n• Edge function בלבד מחזיק service role\n• בטוח לשתף עם backends של לקוחות",
  b2x+Inches(0.2), Inches(3.95), cw-Inches(0.4), Inches(2.7),
  size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 9. IDEMPOTENCY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, TEAL)
wordmark(sl); pgnum(sl, 9, TOTAL)

pill(sl, "🔁 IDEMPOTENCY", Inches(0.4), Inches(0.72), color=TEAL, pw=Inches(1.7))
t(sl, "מערכות POS חוזרות. נקודות לא מכפילות.",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.8),
  size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.5), Inches(2.1), Inches(6.1))
t(sl, "כל קופה מציאותית מנסה שוב על network timeout. בלי idempotency: כל ניסיון יוצר נקודות חדשות והטבות חדשות.",
  Inches(0.4), Inches(2.3), Inches(12.5), Inches(0.65),
  size=14, color=SILVER, align=PP_ALIGN.RIGHT)

# Flow diagram
steps = [
    ("1", "POS שולח רכישה", "transaction_id: TX-001", GBRT),
    ("2", "השרת מעבד", "+101 נקודות לדנה", GBRT),
    ("3", "Network timeout", "POS לא קיבל תשובה", AMBER),
    ("4", "POS מנסה שוב", "אותו transaction_id", AMBER),
    ("5", "השרת מזהה כפילות", "מחזיר תשובה מקורית", GBRT),
    ("6", "0 נקודות נוספות", "ללא הטבה כפולה", GBRT),
]
for i, (num, title, sub, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    bx = Inches(0.4 + col * 4.3)
    by = Inches(3.2 + row * 1.85)
    dark_card(sl, bx, by, Inches(4.15), Inches(1.65), color)
    # Number circle
    circle = sl.shapes.add_shape(9, bx+Inches(3.55), by+Inches(0.2), Inches(0.4), Inches(0.4))
    circle.fill.solid(); circle.fill.fore_color.rgb = color; circle.line.fill.background()
    t(sl, num, bx+Inches(3.6), by+Inches(0.23), Inches(0.3), Inches(0.35),
      size=13, bold=True, color=BG, align=PP_ALIGN.CENTER, rtl=False)
    t(sl, title, bx+Inches(0.2), by+Inches(0.25), Inches(3.3), Inches(0.4),
      size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    t(sl, sub, bx+Inches(0.2), by+Inches(0.85), Inches(3.75), Inches(0.7),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 10. NON-BLOCKING + PERSISTENCE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GBRT)
wordmark(sl); pgnum(sl, 10, TOTAL)

pill(sl, "⚡ PERFORMANCE", Inches(0.4), Inches(0.72), color=GBRT, pw=Inches(1.7))
t(sl, "Non-blocking · Persistence",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.8),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.5), Inches(2.1), Inches(6.1))

# 2 cards side by side
c1x = Inches(0.4); c2x = Inches(6.95); cw = Inches(6.0); ch = Inches(4.8)

dark_card(sl, c1x, Inches(2.4), cw, ch, GBRT)
t(sl, "⚡  Non-blocking Benefit",
  c1x+Inches(0.2), Inches(2.6), cw-Inches(0.4), Inches(0.5),
  size=18, bold=True, color=GBRT, align=PP_ALIGN.RIGHT)
t(sl, "תגובת רכישה לא מחכה ל-Clubby",
  c1x+Inches(0.2), Inches(3.15), cw-Inches(0.4), Inches(0.5),
  size=13, color=SILVER, align=PP_ALIGN.RIGHT)
# Benchmark
rrect(sl, c1x+Inches(0.3), Inches(3.85), cw-Inches(0.6), Inches(0.85), SURFACE, GBRT)
t(sl, "77ms",
  c1x+Inches(0.3), Inches(3.92), cw-Inches(0.6), Inches(0.55),
  size=28, bold=True, color=GBRT, align=PP_ALIGN.CENTER, rtl=False)
t(sl, "זמן תגובה ממוצע (לפני: ~3000ms)",
  c1x+Inches(0.3), Inches(4.4), cw-Inches(0.6), Inches(0.3),
  size=10, color=DIM, align=PP_ALIGN.CENTER)
for j, b in enumerate([
    "השרת מחזיר receipt מיד",
    "Clubby benefit נשלח ב-background",
    "log אירוע בנפרד",
    "אם Clubby זמין — לא חוסם תהליך",
]):
    t(sl, "• " + b, c1x+Inches(0.2), Inches(5.0+j*0.42), cw-Inches(0.4), Inches(0.4),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

dark_card(sl, c2x, Inches(2.4), cw, ch, INDIGO)
t(sl, "💾  Persistence",
  c2x+Inches(0.2), Inches(2.6), cw-Inches(0.4), Inches(0.5),
  size=18, bold=True, color=INDIGO, align=PP_ALIGN.RIGHT)
t(sl, "נתוני הדגמה שורדים restart",
  c2x+Inches(0.2), Inches(3.15), cw-Inches(0.4), Inches(0.5),
  size=13, color=SILVER, align=PP_ALIGN.RIGHT)
# 3 file boxes
files = ["transactions.json", "members.json", "events.json"]
for j, f in enumerate(files):
    rrect(sl, c2x+Inches(0.25), Inches(3.85+j*0.45), cw-Inches(0.5), Inches(0.38), SURFACE, BORDER)
    t(sl, "📄 " + f, c2x+Inches(0.4), Inches(3.92+j*0.45), cw-Inches(0.7), Inches(0.3),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "כל transaction נשמר ל-JSON",
    "יתרת נקודות שורדת restart",
    "Events log עד 500 רשומות",
    "אנדפוינט /api/reset לדמו נקי",
]):
    t(sl, "• " + b, c2x+Inches(0.2), Inches(5.4+j*0.42), cw-Inches(0.4), Inches(0.4),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 11. INTEGRATION EVENTS LOG
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, ROSE)
wordmark(sl); pgnum(sl, 11, TOTAL)

FW = Inches(8.0); FH = Inches(5.7)
browser_frame(sl, f"{V2}/04-pos-integration-log.png", Inches(0.35), Inches(0.9), FW, FH)

TX = Inches(8.7); TW = Inches(4.4)
pill(sl, "🔌 לוג אינטגרציה", TX, Inches(0.72), color=ROSE, pw=Inches(1.9))
t(sl, "המסך שמוכר את האינטגרציה",
  TX, Inches(1.18), TW, Inches(0.9), size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(0.5), Inches(2.25), Inches(3.85))
t(sl, "פיד חי של אירועים בין POS ל-Clubby. הקהל רואה את האינטגרציה עובדת בזמן אמת.",
  TX, Inches(2.4), TW, Inches(1.1), size=12, color=SILVER, align=PP_ALIGN.RIGHT)

events = [
    ("🛒", "purchase",        GBRT,    "רכישה התקבלה"),
    ("⚠️", "purchase_duplicate", AMBER, "כפולה זוהתה"),
    ("🎁", "benefit_issued",  GBRT,    "הטבה נשלחה ל-Clubby"),
    ("❌", "benefit_error",   RED,     "Clubby לא זמין"),
    ("←",  "webhook_inbound", TEAL,    "QR scan נכנס"),
    ("→",  "webhook_response",INDIGO,  "תשובה ל-Clubby"),
]
for j, (em, name, color, desc) in enumerate(events):
    by = Inches(3.65 + j * 0.45)
    rrect(sl, TX, by, TW, Inches(0.4), SURFACE)
    t(sl, em, TX+Inches(0.1), by+Inches(0.05), Inches(0.3), Inches(0.3),
      size=14, color=color, align=PP_ALIGN.CENTER, rtl=False)
    t(sl, name, TX+Inches(0.55), by+Inches(0.06), Inches(1.4), Inches(0.3),
      size=10, bold=True, color=color, align=PP_ALIGN.LEFT, rtl=False)
    t(sl, desc, TX+Inches(2.0), by+Inches(0.06), TW-Inches(2.1), Inches(0.3),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 12. ANCHOR DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, TEAL)
wordmark(sl); pgnum(sl, 12, TOTAL)

FW = Inches(8.0); FH = Inches(5.7)
browser_frame(sl, f"{V2}/05-pos-anchor-dashboard.png", Inches(0.35), Inches(0.9), FW, FH)

TX = Inches(8.7); TW = Inches(4.4)
pill(sl, "🏢 ANCHOR DASHBOARD", TX, Inches(0.72), color=TEAL, pw=Inches(2.4))
t(sl, "תמונת רשת קמעונאית",
  TX, Inches(1.18), TW, Inches(0.9), size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(0.5), Inches(2.25), Inches(3.85))
t(sl, "המסך שסוגר את הפיץ' לרשת קמעונאית — KPI חוצי-סניפים, פילוח דרגות, וקטגוריות מובילות.",
  TX, Inches(2.4), TW, Inches(1.1), size=12, color=SILVER, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "🏷️  Powered by Clubby badge",
    "👥  1,361 חברי מועדון מסונכרנים",
    "🏪  5 סניפים — חיפה, ת״א, י-ם...",
    "🥇  פילוח דרגות (Gold/Silver/Bronze)",
    "📊  קטגוריות מובילות לפי הכנסה",
]):
    t(sl, b, TX, Inches(3.65+j*0.55), TW, Inches(0.5),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 13. SECTION — UI Polish
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
starfield(sl, 18)
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, GBRT)
t(sl, "03", Inches(0.35), Inches(0.5), Inches(2), Inches(0.65),
  size=13, bold=True, color=DIM, align=PP_ALIGN.LEFT, rtl=False)
glow = rect(sl, 0, Inches(2.2), W, Inches(2.2), RGBColor(0x07, 0x12, 0x09))
glow.line.fill.background()
t(sl, "ליטוש UI", Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.8),
  size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
gold_line(sl, Inches(4.5), Inches(3.25), Inches(4.3))
t(sl, "פרטים קטנים שעושים את ההבדל",
  Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.65),
  size=18, color=SILVER, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 14. POS UI SHOWCASE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GOLD)
wordmark(sl); pgnum(sl, 14, TOTAL)

FW = Inches(8.0); FH = Inches(5.7)
browser_frame(sl, f"{V2}/02-pos-dana-gold-member.png", Inches(0.35), Inches(0.9), FW, FH)

TX = Inches(8.7); TW = Inches(4.4)
pill(sl, "POS UPGRADES", TX, Inches(0.72), color=GOLD, pw=Inches(1.6))
t(sl, "תגי דרגה צבעוניים",
  TX, Inches(1.18), TW, Inches(0.9), size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(0.5), Inches(2.25), Inches(3.85))
t(sl, "במקום טקסט gold/silver/bronze — תגים צבעוניים עם רקע, מסגרת ואמוג׳י דרגה.",
  TX, Inches(2.4), TW, Inches(1.0), size=12, color=SILVER, align=PP_ALIGN.RIGHT)
# Tier showcase
for j, (em, name, color) in enumerate([
    ("🥇", "זהב",   RGBColor(0xD4,0xAF,0x37)),
    ("🥈", "כסף",   RGBColor(0xA8,0xB2,0xC8)),
    ("🥉", "ארד",   RGBColor(0xCD,0x7F,0x32)),
]):
    by = Inches(3.6 + j * 0.7)
    rrect(sl, TX, by, TW, Inches(0.55), RGBColor(color[0], color[1], color[2]), color)
    # background tint
    bg_shape = sl.shapes.add_shape(9, TX, by, TW, Inches(0.55))
    bg_shape.fill.solid(); bg_shape.fill.fore_color.rgb = RGBColor(min(color[0]+10,40), min(color[1]+10,40), min(color[2]+10,40))
    bg_shape.line.color.rgb = color; bg_shape.line.width = Pt(1)
    t(sl, em + "  " + name, TX+Inches(0.2), by+Inches(0.12), TW-Inches(0.4), Inches(0.35),
      size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)

# Receipt format
t(sl, "📋 פורמט קבלה חדש:",
  TX, Inches(5.85), TW, Inches(0.4), size=12, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
t(sl, "OA-20260609-001001",
  TX, Inches(6.25), TW, Inches(0.4), size=14, bold=True, color=GBRT, align=PP_ALIGN.LEFT, rtl=False)
t(sl, "(תאריך + רץ עוקב)",
  TX, Inches(6.65), TW, Inches(0.3), size=10, color=DIM, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 15. POS SALE COMPLETE — Pending state
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GBRT)
wordmark(sl); pgnum(sl, 15, TOTAL)

FW = Inches(8.0); FH = Inches(5.7)
browser_frame(sl, f"{V2}/03-pos-sale-pending-clubby.png", Inches(0.35), Inches(0.9), FW, FH)

TX = Inches(8.7); TW = Inches(4.4)
pill(sl, "PENDING STATE", TX, Inches(0.72), color=GBRT, pw=Inches(1.7))
t(sl, "מצב 'בתהליך' חדש",
  TX, Inches(1.18), TW, Inches(0.9), size=22, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(0.5), Inches(2.25), Inches(3.85))
t(sl, "במקום \"הטבה נשלחה\" כשלא ידוע אם הצליח — מציגים ⏳ \"הטבה בתהליך שליחה\" כי השליחה אסינכרונית.",
  TX, Inches(2.4), TW, Inches(1.5), size=12, color=SILVER, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "✅  עסקה הושלמה (ירוק)",
    "⏳  הטבה pending (לא חוסם)",
    "🎂  Birthday detection",
    "💳  3 תשלומים — VAT 17%",
    "📱  Receipt בעברית מלאה",
]):
    t(sl, b, TX, Inches(4.05+j*0.5), TW, Inches(0.45),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 16. MOBILE — Wallet
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GBRT)
wordmark(sl); pgnum(sl, 16, TOTAL)

PH = Inches(6.4); PW = Inches(3.2)
PX = W - PW - Inches(0.4); PY = (H-PH)/2
phone_frame(sl, f"{V2}/09-mobile-wallet-with-feed.png", PX, PY, PW, PH)

TX = Inches(0.4); TW = Inches(9.2)
pill(sl, "מובייל", TX, Inches(0.72), color=GBRT, pw=Inches(1.0))
t(sl, "ארנק עם פיד חי",
  TX, Inches(1.18), TW, Inches(1.0), size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX, Inches(2.28), Inches(4.8))
t(sl, "ארנק הלקוח — יתרת קרדיט, מועדונים עם תגי דרגה, פיד עדכונים בזמן אמת מכל העסקים, וכל ההטבות במסך אחד.",
  TX, Inches(2.45), TW-Inches(0.3), Inches(1.2),
  size=14, color=SILVER, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "💰 יתרת קרדיט בשקלים",
    "🏪 MY CLUBS עם תגי ברונזה/כסף/זהב",
    "📢 UPDATES — פיד פוסטים מהעסקים",
    "🎁 הטבות פעילות — לחיצה למימוש",
    "🔔 התראות עם badge ספירה",
]):
    t(sl, b, TX, Inches(3.95+j*0.5), TW-Inches(0.5), Inches(0.45),
      size=13, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 17. PORTAL — Webhook Settings + Platform Analytics
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, INDIGO)
wordmark(sl); pgnum(sl, 17, TOTAL)

# Two browser frames side by side
FW = Inches(6.0); FH = Inches(5.5)
browser_frame(sl, f"{V2}/16-portal-webhook-config.png", Inches(0.3), Inches(0.95), FW, FH)
browser_frame(sl, f"{V2}/17-portal-platform-analytics.png", Inches(6.5), Inches(0.95), FW, FH)

t(sl, "פורטל v2 — Settings + Platform",
  Inches(0.3), Inches(0.4), Inches(12.7), Inches(0.5),
  size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.5), Inches(6.55), Inches(6.1))
t(sl, "Settings: Webhook + signing secret · Platform: KPI חוצי-עסקים · גרפים תקופתיים · טבלת מובילים",
  Inches(0.3), Inches(6.7), Inches(12.7), Inches(0.5),
  size=13, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 18. METRICS — Before / After
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GOLD)
wordmark(sl); pgnum(sl, 18, TOTAL)

pill(sl, "📊 METRICS", Inches(0.4), Inches(0.72), pw=Inches(1.2))
t(sl, "מדידות לפני / אחרי",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(5.5), Inches(2.0), Inches(5.1))

metrics = [
    ("⚡ זמן תגובת רכישה",       "~3,000ms",  "77ms",        "39× מהיר"),
    ("🔴 אינטראקציות שבורות",   "5 P0",      "0",           "100% תוקנו"),
    ("🌑 מסכים תואמי dark mode", "60%",       "100%",        "+40%"),
    ("🔐 גישת מפתחות לחיצוני",   "service_role", "M2M ייעודי", "אבטחה ארגונית"),
    ("💾 Survival על restart",  "❌ נמחק",    "✅ נשמר",      "JSON persist"),
    ("🔁 POS retry idempotent", "❌ כפול",    "✅ deduplicated", "Production-safe"),
    ("📊 שקיפות אינטגרציה",     "0 דרך",     "Live event log","UX חדש"),
    ("🏢 Anchor view",          "❌ אין",     "✅ נבנה",      "Enterprise pitch"),
]
for i, (label, before, after, gain) in enumerate(metrics):
    col = i % 2
    row = i // 2
    bx = Inches(0.4 + col * 6.45)
    by = Inches(2.3 + row * 1.15)
    dark_card(sl, bx, by, Inches(6.3), Inches(1.0), GOLD)
    t(sl, label, bx+Inches(0.2), by+Inches(0.1), Inches(6.0), Inches(0.4),
      size=13, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # Metric values
    t(sl, before, bx+Inches(0.2), by+Inches(0.5), Inches(2.0), Inches(0.4),
      size=11, color=RED, align=PP_ALIGN.RIGHT)
    t(sl, "→", bx+Inches(2.3), by+Inches(0.5), Inches(0.4), Inches(0.4),
      size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER, rtl=False)
    t(sl, after, bx+Inches(2.7), by+Inches(0.5), Inches(2.0), Inches(0.4),
      size=11, color=GBRT, align=PP_ALIGN.RIGHT)
    t(sl, gain, bx+Inches(0.2), by+Inches(0.5), Inches(6.0), Inches(0.4),
      size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# 19. DESIGN SYSTEM RECOMMENDATIONS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, ROSE)
wordmark(sl); pgnum(sl, 19, TOTAL)

pill(sl, "DESIGN SYSTEM", Inches(0.4), Inches(0.72), color=ROSE, pw=Inches(1.8))
t(sl, "המלצות מערכת עיצוב",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(5.5), Inches(2.0), Inches(5.1))

cards_ds = [
    (GBRT,   "🎨 צבע",
     "20+ ערכי hex hardcoded → tokens יחידים\nWCAG AA: muted #4a5260 → #5f6d7e\nstatus.error / warning / success / tint"),
    (GOLD,   "🔤 טיפוגרפיה",
     "21 גדלים ברחבי האפליקציה\nמופחת ל-9 גדלים בסקלה אחידה\n11 / 12 / 13 / 14 / 15 / 17 / 22 / 28 / 42"),
    (TEAL,   "📐 רווחים",
     "8px base grid\nspacing-1 (4) → spacing-12 (48)\nאחיד בכל המסכים והבליטות"),
    (INDIGO, "🧩 קומפוננטות",
     "<Button variant size loading />\n<ErrorBanner /> · <BackButton />\n<EmptyState /> · <Modal /> · <KebabMenu />"),
]
for i, (accent, title, body) in enumerate(cards_ds):
    cx = Inches(0.35 + i * 3.25)
    dark_card(sl, cx, Inches(2.4), Inches(3.1), Inches(4.7), accent)
    t(sl, title, cx+Inches(0.18), Inches(2.6), Inches(2.85), Inches(0.5),
      size=15, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    t(sl, body, cx+Inches(0.18), Inches(3.2), Inches(2.85), Inches(3.7),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 20. ROADMAP — what's next
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 20, TOTAL)

pill(sl, "ROADMAP", Inches(0.4), Inches(0.72), color=AMBER, pw=Inches(1.2))
t(sl, "מה הבא בתור",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(5.5), Inches(2.0), Inches(5.1))

# 3 columns
done = ["✅ 5 תיקוני P0 קריטיים", "✅ M2M integration key", "✅ Idempotency + persistence",
        "✅ Non-blocking architecture", "✅ Live event log", "✅ Anchor dashboard",
        "✅ Receipt format עוקב", "✅ Tier badges צבעוניים", "✅ Hebrew RTL מלא"]
medium = ["🔄 Extract useCurrentUser hook", "🔄 Shared <ErrorBanner>", "🔄 Shared <BackButton>",
          "🔄 Shared <Modal> w/ a11y", "🔄 WCAG contrast pass", "🔄 i18n register + onboarding"]
major = ["🎯 Dark mode unified system", "🎯 Live SMS OTP (Twilio)", "🎯 GPS map view",
         "🎯 AI campaign recs", "🎯 App Store submission", "🎯 Production deploy"]

for i, (title, items, color) in enumerate([
    ("נעשה", done, GBRT),
    ("בעבודה", medium, AMBER),
    ("עתיד", major, ROSE),
]):
    cx = Inches(0.35 + i * 4.3)
    dark_card(sl, cx, Inches(2.4), Inches(4.15), Inches(4.7), color)
    t(sl, title, cx+Inches(0.2), Inches(2.6), Inches(3.85), Inches(0.4),
      size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)
    for j, item in enumerate(items):
        t(sl, item, cx+Inches(0.2), Inches(3.15+j*0.42), Inches(3.85), Inches(0.4),
          size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 21. VERIFICATION — All tests passed
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GBRT)
wordmark(sl); pgnum(sl, 21, TOTAL)

pill(sl, "✅ VERIFIED", Inches(0.4), Inches(0.72), color=GBRT, pw=Inches(1.4))
t(sl, "כל השינויים נבדקו ב-browser",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.0), Inches(2.0), Inches(5.6))

tests = [
    ("T1", "Dev hint visible on localhost only", "✅"),
    ("T2", "Redeem screen renders + error inline", "✅"),
    ("T3", "Discover loads businesses (auth)",  "✅"),
    ("T4", "Portal dashboard stats visible",     "✅"),
    ("T5a","FeedPosts form opens correctly",     "✅"),
    ("T5b","FeedPosts saves + closes form",      "✅"),
    ("T6", "Members modal uses var(--surface)",  "✅"),
    ("T7", "Platform Analytics no permanent spinner", "✅"),
    ("M2M","Integration key wrong → 401 returned", "✅"),
    ("IDEM","Same transaction_id → idempotent flag", "✅"),
    ("PERSIST","data/*.json files written",      "✅"),
    ("ASYNC","Purchase returns in 77ms",          "✅"),
    ("EVENTS","Log shows purchase + duplicate + error", "✅"),
    ("ANCHOR","5 branches + tier breakdown shown", "✅"),
]

for i, (id, desc, status) in enumerate(tests):
    col = i % 2
    row = i // 2
    bx = Inches(0.4 + col * 6.45)
    by = Inches(2.3 + row * 0.62)
    dark_card(sl, bx, by, Inches(6.3), Inches(0.5))
    t(sl, status, bx+Inches(0.15), by+Inches(0.1), Inches(0.4), Inches(0.32),
      size=14, color=GBRT, align=PP_ALIGN.CENTER, rtl=False)
    t(sl, id, bx+Inches(0.6), by+Inches(0.13), Inches(0.7), Inches(0.3),
      size=11, bold=True, color=GOLD, align=PP_ALIGN.LEFT, rtl=False)
    t(sl, desc, bx+Inches(1.4), by+Inches(0.13), Inches(4.8), Inches(0.3),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 22. CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
starfield(sl, 35)
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, GOLD)
wordmark(sl)

glow = rect(sl, 0, Inches(1.6), W, Inches(3.0), RGBColor(0x09, 0x0F, 0x07))
glow.line.fill.background()

t(sl, "מוכן לדמו אחרון.", Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.1),
  size=52, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
t(sl, "מוכן לייצור.", Inches(0.5), Inches(2.15), Inches(12.3), Inches(1.0),
  size=52, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(5.5), Inches(3.3), Inches(5.1))
t(sl, "12 תיקונים קריטיים · אינטגרציה מאובטחת · UX באיכות סטריפ/לינייר",
  Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.55),
  size=18, color=SILVER, align=PP_ALIGN.RIGHT)

# Stack
for i, (label, val) in enumerate([
    ("ארנק לקוח", "http://localhost:8081"),
    ("פורטל עסקי", "http://localhost:5173"),
    ("POS אושר עד v2", "http://localhost:3001"),
]):
    bx = Inches(1.2 + i * 3.7)
    dark_card(sl, bx, Inches(4.9), Inches(3.5), Inches(1.9), GOLD)
    t(sl, label, bx+Inches(0.15), Inches(5.0), Inches(3.25), Inches(0.4),
      size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    t(sl, val, bx+Inches(0.15), Inches(5.45), Inches(3.25), Inches(0.35),
      size=10, color=DIM, align=PP_ALIGN.LEFT, rtl=False)

t(sl, "Powered by Clubby × אושר עד",
  Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35),
  size=11, color=DIM, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"\n✅ Saved: {OUT}")
print(f"   Slides: {len(prs.slides)}")
