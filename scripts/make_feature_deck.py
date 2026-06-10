"""
Clubby — Full Feature Deck (English, with real demo screenshots)
Covers: customer app, business portal, Osher-Ad integration
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

BASE     = "/Users/I560647/Library/CloudStorage/OneDrive-SAPSE/Desktop/finance/Clubby"
D        = f"{BASE}/screenshots/demo"
S        = f"{BASE}/screenshots/scenario"
WORDMARK = f"{BASE}/logo/wordmark.png"
OUT      = f"{BASE}/slides/Clubby_Feature_Deck_2026.pptx"
os.makedirs(f"{BASE}/slides", exist_ok=True)

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x09, 0x0E, 0x1A)
GREEN  = RGBColor(0x1A, 0x7A, 0x4A)
GBRT   = RGBColor(0x2E, 0xCC, 0x71)
INDIGO = RGBColor(0x5B, 0x4C, 0xF5)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
PINK   = RGBColor(0xEC, 0x48, 0x99)
TEAL   = RGBColor(0x06, 0xB6, 0xD4)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFF    = RGBColor(0xF8, 0xFA, 0xFC)
GRAY   = RGBColor(0x64, 0x74, 0x8B)
LGRAY  = RGBColor(0x94, 0xA3, 0xB8)
DARK   = RGBColor(0x1E, 0x29, 0x3B)
DKGRAY = RGBColor(0x2D, 0x3A, 0x4A)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def fill(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def t(slide, text, x, y, w, h, size=18, bold=False, color=DARK,
      align=PP_ALIGN.LEFT, italic=False):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color

def rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def img(slide, path, x, y, w, h=None):
    if not path or not os.path.exists(path): return False
    if h: slide.shapes.add_picture(path, x, y, w, h)
    else: slide.shapes.add_picture(path, x, y, w)
    return True

def phone_frame(slide, path, x, y, fw, fh):
    if not path or not os.path.exists(path): return
    frm = slide.shapes.add_shape(9, x, y, fw, fh)
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0x10, 0x14, 0x1E)
    frm.line.color.rgb = RGBColor(0x2A, 0x30, 0x4A); frm.line.width = Pt(2)
    pad = Inches(0.1)
    slide.shapes.add_picture(path, x+pad, y+pad, fw-pad*2, fh-pad*2)

def browser_frame(slide, path, x, y, fw, fh):
    if not path or not os.path.exists(path): return
    chrome_h = Inches(0.22)
    frm = slide.shapes.add_shape(1, x, y, fw, fh)
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    frm.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1); frm.line.width = Pt(1)
    chrome = slide.shapes.add_shape(1, x, y, fw, chrome_h)
    chrome.fill.solid(); chrome.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    chrome.line.fill.background()
    # dots
    for di, dc in enumerate([RGBColor(0xFC,0x5C,0x5C), RGBColor(0xFD,0xBC,0x40), RGBColor(0x27,0xC9,0x3F)]):
        d = slide.shapes.add_shape(9, x+Inches(0.08+di*0.17), y+Inches(0.07), Inches(0.09), Inches(0.09))
        d.fill.solid(); d.fill.fore_color.rgb = dc; d.line.fill.background()
    slide.shapes.add_picture(path, x, y+chrome_h, fw, fh-chrome_h)

def pill(slide, label, accent, x, y, pw=Inches(1.4), ph=Inches(0.3)):
    p = slide.shapes.add_shape(9, x, y, pw, ph)
    p.fill.solid(); p.fill.fore_color.rgb = accent
    p.line.fill.background()
    t(slide, label, x+Inches(0.06), y+Inches(0.04), pw-Inches(0.1), ph-Inches(0.06),
      size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def pgnum(slide, n, total):
    t(slide, f"{n}  /  {total}", Inches(0.3), H-Inches(0.38), Inches(1.5), Inches(0.3),
      size=9, color=LGRAY, align=PP_ALIGN.LEFT)

def wordmark(slide, light=False):
    if os.path.exists(WORDMARK):
        slide.shapes.add_picture(WORDMARK, Inches(0.28), Inches(0.18), Inches(2.2), Inches(0.6))
    else:
        t(slide, "clubby.", Inches(0.28), Inches(0.18), Inches(2.5), Inches(0.55),
          size=20, bold=True, color=GBRT if not light else WHITE, align=PP_ALIGN.LEFT)

TOTAL = 26  # total slides

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, NAVY)
rect(sl, 0, 0, Inches(0.07), H, GBRT)
wordmark(sl)

# Big hero text
t(sl, "Clubby.", Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.6),
  size=72, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
t(sl, "The loyalty platform that plugs into your existing backend.",
  Inches(0.5), Inches(3.0), Inches(12.3), Inches(0.9),
  size=24, color=LGRAY, align=PP_ALIGN.RIGHT)

# Bottom stat band
rect(sl, 0, Inches(4.6), W, Inches(2.2), RGBColor(0x0F,0x17,0x2A))
stats = [
    ("Customer Wallet App", "iOS · Android · Web"),
    ("Business Portal", "Full CRM + Analytics"),
    ("Backend Integration", "Webhooks · REST API"),
    ("Platform Analytics", "Cross-business insights"),
]
for i, (a, b) in enumerate(stats):
    x = Inches(0.3 + i * 3.25)
    rect(sl, x, Inches(4.75), Inches(3.0), Inches(1.75), RGBColor(0x14,0x20,0x33))
    t(sl, a, x+Inches(0.15), Inches(4.92), Inches(2.7), Inches(0.5),
      size=14, bold=True, color=GBRT, align=PP_ALIGN.LEFT)
    t(sl, b, x+Inches(0.15), Inches(5.45), Inches(2.7), Inches(0.4),
      size=11, color=LGRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS CLUBBY (overview)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, GREEN)
wordmark(sl)
pgnum(sl, 2, TOTAL)

t(sl, "One platform. Three products.", Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.8),
  size=36, bold=True, color=DARK, align=PP_ALIGN.LEFT)
t(sl, "Clubby connects loyalty programs, store backends, and customer wallets in a single integrated platform.",
  Inches(0.4), Inches(1.95), Inches(12.5), Inches(0.7),
  size=16, color=GRAY, align=PP_ALIGN.LEFT)
rect(sl, Inches(0.4), Inches(2.75), Inches(4.0), Inches(0.04), GREEN)

cards = [
    (GREEN,  "📱", "Customer Wallet",
     ["Phone OTP sign-in, no password",
      "Loyalty benefits, stamp cards, points",
      "QR scan to join any club",
      "Real-time feed from businesses",
      "Dark mode, Hebrew RTL, 5 languages"]),
    (INDIGO, "🏢", "Business Portal",
     ["Full CRM — members, tiers, points",
      "Promotions, stamp cards, stories",
      "Push notifications (segmented)",
      "Analytics dashboard + CSV export",
      "Feed posts to member wallets"]),
    (AMBER,  "🔌", "Backend Integration",
     ["Webhook: QR scan → your backend",
      "REST API: push benefits from POS",
      "HMAC-SHA256 signature verification",
      "Tier-aware benefit logic",
      "Works with any existing backend"]),
]
for i, (accent, emoji, title, bullets) in enumerate(cards):
    cx = Inches(0.4 + i * 4.3)
    card = sl.shapes.add_shape(1, cx, Inches(3.0), Inches(4.1), Inches(4.0))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2,0xE8,0xF0); card.line.width = Pt(1)
    rect(sl, cx, Inches(3.0), Inches(4.1), Inches(0.06), accent)
    t(sl, emoji + "  " + title, cx+Inches(0.18), Inches(3.22), Inches(3.7), Inches(0.5),
      size=16, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    for j, b in enumerate(bullets):
        dot = sl.shapes.add_shape(9, cx+Inches(0.2), Inches(3.92+j*0.5)+Inches(0.06),
                                   Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        t(sl, b, cx+Inches(0.38), Inches(3.9+j*0.5), Inches(3.6), Inches(0.45),
          size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SECTION DIVIDER: Customer App
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, GREEN)
rect(sl, 0, 0, Inches(0.07), H, WHITE)
t(sl, "01", Inches(0.3), Inches(0.5), Inches(2), Inches(0.7),
  size=14, bold=True, color=RGBColor(0xFF,0xFF,0xFF,), align=PP_ALIGN.LEFT)
t(sl, "Customer Wallet App", Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5),
  size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t(sl, "Sign in · Join clubs · Earn benefits · Redeem · Discover",
  Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.7),
  size=20, color=RGBColor(0xBB,0xF7,0xD0), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Sign-in + Verify
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, GREEN)
wordmark(sl)
pgnum(sl, 4, TOTAL)

pill(sl, "AUTH", GREEN, Inches(0.4), Inches(0.75))
t(sl, "Phone OTP — no password, no friction",
  Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.8),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, Inches(0.4), Inches(2.15), Inches(3.5), Inches(0.04), GREEN)
t(sl, "Customers sign in with their phone number. One-tap OTP — no accounts, no passwords. Unknown numbers show a clear inline error directing them to scan a store QR first.",
  Inches(0.4), Inches(2.3), Inches(5.8), Inches(1.4),
  size=14, color=GRAY, align=PP_ALIGN.LEFT)
features = [
    "✓  Phone lookup — only known members proceed",
    "✓  Test OTP 000000 in local dev",
    "✓  Local dev skips OTP API — goes straight to verify",
    "✓  Unknown phone → inline red error (no Alert popups)",
    "✓  OTP verified → wallet or register",
]
for j, f in enumerate(features):
    t(sl, f, Inches(0.4), Inches(3.85+j*0.52), Inches(5.8), Inches(0.45),
      size=13, color=DARK, align=PP_ALIGN.LEFT)

# Two phones side by side
PH = Inches(5.8); PW = Inches(2.9)
phone_frame(sl, f"{D}/11-mobile-signin-screen.png", Inches(6.5), Inches(0.75), PW, PH)
phone_frame(sl, f"{D}/12-mobile-verify-screen.png", Inches(9.7), Inches(0.75), PW, PH)
t(sl, "Sign-in", Inches(6.8), Inches(6.6), Inches(2.5), Inches(0.3),
  size=10, color=LGRAY, align=PP_ALIGN.CENTER)
t(sl, "OTP Verify", Inches(10.0), Inches(6.6), Inches(2.5), Inches(0.3),
  size=10, color=LGRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Customer Wallet
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, GREEN)
wordmark(sl)
pgnum(sl, 5, TOTAL)

PH = Inches(6.5); PW = Inches(3.25)
phone_frame(sl, f"{S}/P2-dana-wallet.png", Inches(0.45), Inches(0.5), PW, PH)

TX = Inches(4.2)
pill(sl, "WALLET", GREEN, TX, Inches(0.65))
t(sl, "Everything in one place",
  TX, Inches(1.15), Inches(8.8), Inches(0.7),
  size=30, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(5.0), Inches(0.04), GREEN)
t(sl, "The customer's loyalty hub — balance, clubs, feed updates, benefits and stamp progress all in a single dark-themed app.",
  TX, Inches(2.15), Inches(8.8), Inches(1.0),
  size=14, color=GRAY, align=PP_ALIGN.LEFT)
items = [
    ("🏦", "Balance card", "Total credit in ILS (₪), active count"),
    ("🏪", "My Clubs",     "Each joined business with tier badge (Bronze/Silver/Gold)"),
    ("📢", "Updates feed", "Live posts from businesses — offers, announcements"),
    ("🎁", "Benefits list","All/Balance/Discounts/Free Items tabs, redeem in one tap"),
    ("🔔", "Notifications","Bell icon with unread badge, full notification history"),
]
for j, (em, title, sub) in enumerate(items):
    by = Inches(3.4 + j * 0.65)
    t(sl, em, TX, by, Inches(0.4), Inches(0.55), size=16, align=PP_ALIGN.LEFT)
    t(sl, title, TX+Inches(0.45), by, Inches(2.2), Inches(0.3),
      size=13, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    t(sl, sub, TX+Inches(0.45), by+Inches(0.28), Inches(8.4), Inches(0.32),
      size=11, color=GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Discover + Store Profile
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, TEAL)
wordmark(sl)
pgnum(sl, 6, TOTAL)

PH = Inches(5.9); PW = Inches(2.95)
phone_frame(sl, f"{D}/10-mobile-dana-discover.png", Inches(0.3), Inches(0.75), PW, PH)
phone_frame(sl, f"{S}/P4-dana-store-member.png",    Inches(3.55), Inches(0.75), PW, PH)

TX = Inches(7.0)
pill(sl, "DISCOVER", TEAL, TX, Inches(0.65))
t(sl, "Find & join local businesses",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), TEAL)
t(sl, "Discover nearby stores with colorful category squares, search, and Open Now filter. Tap any business to see its full profile, promotions, hours, and join instantly.",
  TX, Inches(2.15), Inches(6.0), Inches(1.2),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
items = [
    "🌐  7 category squares (Food, Clothing, Tech…)",
    "📍  GPS — sorts by distance when location granted",
    "🔍  Debounced search across all businesses",
    "🟢  Open Now filter (based on opening hours)",
    "✓   Joined indicator — no double-enroll",
    "🏪  Store profile: cover, logo, offers, hours",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(3.6 + j*0.52), Inches(5.9), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — QR Scan + Redeem
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, AMBER)
wordmark(sl)
pgnum(sl, 7, TOTAL)

PH = Inches(5.9); PW = Inches(2.95)
phone_frame(sl, f"{S}/P4-dana-wallet-benefit.png", Inches(0.3), Inches(0.75), PW, PH)
phone_frame(sl, f"{D}/14-mobile-redeem-screen.png", Inches(3.55), Inches(0.75), PW, PH)

TX = Inches(7.0)
pill(sl, "QR SCAN & REDEEM", AMBER, TX, Inches(0.65))
t(sl, "Scan to earn. Tap to redeem.",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), AMBER)
t(sl, "Customers scan the store QR with the in-app camera or web fallback. Benefits land in the wallet instantly. Redeem with a single tap — show the screen to the cashier.",
  TX, Inches(2.15), Inches(6.0), Inches(1.2),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
items = [
    "📷  Native camera + jsQR web fallback",
    "⚡  Benefit in wallet within 2 seconds",
    "🏷️  Discount / Credit / Free Item benefit types",
    "✅  Verifiable redemption screen",
    "🔒  HMAC-signed webhook — store controls the benefit",
    "🔄  Idempotent — no duplicate benefits on retry",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(3.6 + j*0.52), Inches(5.9), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Stamp Cards + Tiers + Points
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, PINK)
wordmark(sl)
pgnum(sl, 8, TOTAL)

pill(sl, "LOYALTY MECHANICS", PINK, Inches(0.4), Inches(0.75))
t(sl, "Stamp cards · Tiers · Points",
  Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.75),
  size=32, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, Inches(0.4), Inches(2.1), Inches(5.0), Inches(0.04), PINK)
t(sl, "Three retention mechanics, each independently configurable per business.",
  Inches(0.4), Inches(2.25), Inches(12.5), Inches(0.6),
  size=15, color=GRAY, align=PP_ALIGN.LEFT)

mech = [
    (AMBER, "🥇 Stamp Cards",
     ["Business defines required stamps + reward title",
      "Customer sees progress ring in wallet",
      "Auto-completes → benefit issued instantly",
      "Staff stamps via portal (customer UUID paste)"]),
    (INDIGO, "🏅 Loyalty Tiers",
     ["Owner defines Bronze/Silver/Gold thresholds",
      "Based on total lifetime stamps",
      "Tier badge shown in wallet My Clubs section",
      "Next tier progress shown in store profile"]),
    (TEAL, "⭐ Points Programs",
     ["Points earned per QR scan (configurable rate)",
      "Per-business balance shown in wallet",
      "Rewards catalog: free item, discount, credit",
      "Leaderboard in portal for top spenders"]),
]
for i, (accent, title, bullets) in enumerate(mech):
    cx = Inches(0.4 + i * 4.3)
    card = sl.shapes.add_shape(1, cx, Inches(3.1), Inches(4.1), Inches(3.9))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2,0xE8,0xF0); card.line.width = Pt(1)
    rect(sl, cx, Inches(3.1), Inches(4.1), Inches(0.06), accent)
    t(sl, title, cx+Inches(0.18), Inches(3.3), Inches(3.7), Inches(0.45),
      size=15, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    for j, b in enumerate(bullets):
        dot = sl.shapes.add_shape(9, cx+Inches(0.2), Inches(3.95+j*0.52)+Inches(0.05),
                                   Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        t(sl, b, cx+Inches(0.38), Inches(3.93+j*0.52), Inches(3.6), Inches(0.45),
          size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Manual Coupon + Feed
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, INDIGO)
wordmark(sl)
pgnum(sl, 9, TOTAL)

PH = Inches(5.9); PW = Inches(2.95)
phone_frame(sl, f"{D}/13-mobile-manual-coupon.png", Inches(0.3), Inches(0.75), PW, PH)
phone_frame(sl, f"{S}/P6-dana-feed.png",             Inches(3.55), Inches(0.75), PW, PH)

TX = Inches(7.0)
pill(sl, "EXTRA FEATURES", INDIGO, TX, Inches(0.65))
t(sl, "Manual coupons & live feed",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), INDIGO)

t(sl, "Manual Coupon Entry",
  TX, Inches(2.2), Inches(6.0), Inches(0.4),
  size=14, bold=True, color=DARK, align=PP_ALIGN.LEFT)
manual_items = [
    "Store name, title, type (Credit/Discount/Free item)",
    "Type selector squares change button accent color",
    "Inline validation — no Alert() popups",
    "Marked as 'Manual' badge — not store-verified",
]
for j, b in enumerate(manual_items):
    t(sl, "• " + b, TX, Inches(2.7+j*0.42), Inches(6.0), Inches(0.38),
      size=12, color=GRAY, align=PP_ALIGN.LEFT)

t(sl, "Updates Feed (WhatsApp-style)",
  TX, Inches(4.55), Inches(6.0), Inches(0.4),
  size=14, bold=True, color=DARK, align=PP_ALIGN.LEFT)
feed_items = [
    "Posts appear in wallet from all joined businesses",
    "Types: Offer, Announcement, Promotion, Story",
    "CTA button with custom text",
    "Expires automatically on set date",
    "Owner manages posts in portal → Feed Posts",
]
for j, b in enumerate(feed_items):
    t(sl, "• " + b, TX, Inches(5.05+j*0.42), Inches(6.0), Inches(0.38),
      size=12, color=GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SECTION DIVIDER: Business Portal
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, INDIGO)
rect(sl, 0, 0, Inches(0.07), H, WHITE)
t(sl, "02", Inches(0.3), Inches(0.5), Inches(2), Inches(0.7),
  size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
t(sl, "Business Portal", Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5),
  size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t(sl, "Dashboard · Members · Promotions · Stories · Feed · Analytics",
  Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.7),
  size=20, color=RGBColor(0xC7,0xD2,0xFE), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Portal Dashboard
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, INDIGO)
wordmark(sl)
pgnum(sl, 11, TOTAL)

FW = Inches(6.2); FH = Inches(5.8)
browser_frame(sl, f"{D}/15-portal-dashboard.png", Inches(0.3), Inches(0.75), FW, FH)

TX = Inches(6.9)
pill(sl, "PORTAL", INDIGO, TX, Inches(0.65))
t(sl, "Business dashboard",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), INDIGO)
t(sl, "Real-time stats + 30-day trend charts for every business. Owner sees members, active promotions, benefits issued, and redemption rate at a glance.",
  TX, Inches(2.15), Inches(6.0), Inches(1.1),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
items = [
    "📊  4 stat cards — members, promos, benefits, rate",
    "📈  New members chart (last 30 days)",
    "📉  Daily redemptions chart",
    "⬇️   Export Analytics → full CSV download",
    "🔀  Multi-business switcher in sidebar",
    "🌐  Language toggle (English / עברית)",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(3.5+j*0.52), Inches(6.0), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Members CRM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, GREEN)
wordmark(sl)
pgnum(sl, 12, TOTAL)

FW = Inches(6.2); FH = Inches(5.8)
browser_frame(sl, f"{D}/16-portal-members.png", Inches(0.3), Inches(0.75), FW, FH)

TX = Inches(6.9)
pill(sl, "MEMBERS", GREEN, TX, Inches(0.65))
t(sl, "Full member CRM",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), GREEN)
t(sl, "Complete view of every loyalty club member — name, phone, join date, and one-click detail modal showing their full benefit history.",
  TX, Inches(2.15), Inches(6.0), Inches(1.1),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
items = [
    "👥  Name, phone, join date for every member",
    "🔍  Click View → benefit history modal",
    "📣  Notify members — segmented push notifications",
    "🎯  Audience: All / New / Active / Lapsed",
    "⬇️   Export CSV — member list with all fields",
    "🎁  Stories management with 24h expiry",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(3.5+j*0.52), Inches(6.0), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Promotions + Stamp Cards Portal
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, AMBER)
wordmark(sl)
pgnum(sl, 13, TOTAL)

FW = Inches(6.2); FH = Inches(5.8)
browser_frame(sl, f"{S}/P1-owner-stamp-cards.png", Inches(0.3), Inches(0.75), FW, FH)

TX = Inches(6.9)
pill(sl, "PROMOTIONS", AMBER, TX, Inches(0.65))
t(sl, "Promotions, stamps & tiers",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), AMBER)
t(sl, "Complete loyalty mechanics management — create welcome offers, stamp card programs, tier thresholds, points rewards, and multi-location QR codes.",
  TX, Inches(2.15), Inches(6.0), Inches(1.1),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
items = [
    "🎁  Promotions CRUD — edit, duplicate, delete",
    "🥇  Stamp cards — required stamps, reward title",
    "🏅  Tier definitions — Bronze/Silver/Gold thresholds",
    "⭐  Points programs — earn rate + rewards catalog",
    "📱  Multi-location QR — separate QR per branch",
    "🔗  Webhook — connect your existing backend",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(3.5+j*0.52), Inches(6.0), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Feed Posts + Stories
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, PINK)
wordmark(sl)
pgnum(sl, 14, TOTAL)

FW = Inches(6.2); FH = Inches(5.8)
browser_frame(sl, f"{D}/19-portal-feed-posts.png", Inches(0.3), Inches(0.75), FW, FH)

TX = Inches(6.9)
pill(sl, "ENGAGEMENT", PINK, TX, Inches(0.65))
t(sl, "Feed posts & stories",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), PINK)
t(sl, "Businesses post directly to their members' wallet feed. Stories appear as Instagram-style circles. Both are managed from the portal with full CRUD.",
  TX, Inches(2.15), Inches(6.0), Inches(1.1),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
items = [
    "📣  Feed posts: Offer / Announcement / Promotion",
    "🔘  CTA button with custom text",
    "📅  Auto-expire on set date (RLS-filtered)",
    "⏸️   Pause / Activate / Delete without reload",
    "📸  Stories: full-screen, animated progress bars",
    "👁️   Members see only posts from joined clubs",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(3.5+j*0.52), Inches(6.0), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — SECTION DIVIDER: Backend Integration
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, AMBER)
rect(sl, 0, 0, Inches(0.07), H, WHITE)
t(sl, "03", Inches(0.3), Inches(0.5), Inches(2), Inches(0.7),
  size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
t(sl, "Backend Integration", Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.5),
  size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t(sl, "Connect your existing POS / ERP to Clubby in under an hour",
  Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.7),
  size=20, color=RGBColor(0xFE,0xF3,0xC7), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Webhook Settings (Portal)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, AMBER)
wordmark(sl)
pgnum(sl, 16, TOTAL)

FW = Inches(6.0); FH = Inches(5.6)
browser_frame(sl, f"{D}/17-portal-settings-webhook.png", Inches(0.3), Inches(0.85), FW, FH)

TX = Inches(6.7)
pill(sl, "WEBHOOK SETUP", AMBER, TX, Inches(0.65))
t(sl, "3-field setup. Done.",
  TX, Inches(1.15), Inches(6.2), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.5), Inches(0.04), AMBER)
t(sl, "Enter your endpoint URL and an optional signing secret. Clubby signs every request with HMAC-SHA256 so you can verify it's genuine.",
  TX, Inches(2.15), Inches(6.1), Inches(1.0),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)

# Flow diagram inline
steps = [
    ("1", "Customer scans QR"),
    ("2", "Clubby calls your endpoint"),
    ("3", "Your backend returns benefit"),
    ("4", "Benefit appears in wallet"),
]
for j, (n, label) in enumerate(steps):
    bx = TX + j * Inches(1.55)
    circle = sl.shapes.add_shape(9, bx, Inches(3.45), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = AMBER
    circle.line.fill.background()
    t(sl, n, bx+Inches(0.12), Inches(3.5), Inches(0.3), Inches(0.4),
      size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t(sl, label, bx-Inches(0.1), Inches(4.05), Inches(1.7), Inches(0.5),
      size=10, color=DARK, align=PP_ALIGN.CENTER)
    if j < 3:
        rect(sl, bx+Inches(0.55), Inches(3.66), Inches(0.95), Inches(0.04), LGRAY)

items = [
    "🔐  HMAC-SHA256 — X-Clubby-Signature header",
    "⏱️   5-second timeout with automatic fallback",
    "📋  Fallback: active promotions → no benefit",
    "📄  Webhook Integration Guide built into portal",
    "🧪  ngrok / localtunnel for local dev testing",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(4.75+j*0.44), Inches(6.1), Inches(0.4),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Osher-Ad POS Demo (lookup)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, AMBER)
wordmark(sl)
pgnum(sl, 17, TOTAL)

FW = Inches(7.8); FH = Inches(5.8)
browser_frame(sl, f"{D}/02-pos-member-lookup-dana.png", Inches(0.3), Inches(0.85), FW, FH)

# Label overlay
pill(sl, "LIVE DEMO — OSHER-AD POS", AMBER, Inches(0.45), Inches(6.75))

TX = Inches(8.5)
t(sl, "Osher-Ad Mock POS", TX, Inches(1.0), Inches(4.6), Inches(0.6),
  size=20, bold=True, color=DARK, align=PP_ALIGN.LEFT)
t(sl, "Member: Dana Levi\nTier: 🥇 Gold\nPoints: 2,840\nSpent YTD: ₪5,680",
  TX, Inches(1.75), Inches(4.5), Inches(1.8),
  size=13, color=DARK, align=PP_ALIGN.LEFT)
t(sl, "15 Israeli products\nwith Hebrew names\nand EAN-13 barcodes",
  TX, Inches(3.7), Inches(4.5), Inches(1.5),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
t(sl, "Quick-select:\nDana · Yossi · Maya\nOren · Tamar",
  TX, Inches(5.2), Inches(4.5), Inches(1.5),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Osher-Ad POS Sale Complete
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, AMBER)
wordmark(sl)
pgnum(sl, 18, TOTAL)

FW = Inches(7.8); FH = Inches(5.8)
browser_frame(sl, f"{D}/04-pos-sale-complete-clubby-benefit.png", Inches(0.3), Inches(0.85), FW, FH)
pill(sl, "LIVE DEMO — SALE COMPLETE + CLUBBY BENEFIT", AMBER, Inches(0.45), Inches(6.75))

TX = Inches(8.5)
pill(sl, "RESULT", GREEN, TX, Inches(1.0))
t(sl, "POS sale fired.\nClubby benefit issued.",
  TX, Inches(1.4), Inches(4.6), Inches(0.8),
  size=18, bold=True, color=DARK, align=PP_ALIGN.LEFT)
for b in [
    "🎂  Birthday month → free item benefit",
    "🥇  Gold ×2 points multiplier",
    "📋  Hebrew receipt with VAT",
    "💳  3 installments (תשלומים)",
    "✅  Green banner: הטבה נשלחה ל-Clubby!",
]:
    t(sl, b, TX, Inches(2.35 + ["🎂","🥇","📋","💳","✅"].index(b[0])*0.52),
      Inches(4.6), Inches(0.45), size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — POS Dashboard + API Docs
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, AMBER)
wordmark(sl)
pgnum(sl, 19, TOTAL)

FW = Inches(6.2); FH = Inches(3.2)
browser_frame(sl, f"{D}/07-pos-dashboard.png",  Inches(0.3), Inches(0.8), FW, FH)
browser_frame(sl, f"{D}/08-pos-api-docs.png",   Inches(0.3), Inches(4.1), FW, Inches(3.0))

TX = Inches(6.9)
pill(sl, "INTEGRATION", AMBER, TX, Inches(0.65))
t(sl, "POS dashboard & REST API",
  TX, Inches(1.15), Inches(6.1), Inches(0.7),
  size=26, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.0), Inches(0.04), AMBER)

endpoints = [
    ("GET",  "/api/members?phone=…",     "Member lookup with tier + balance"),
    ("POST", "/api/purchase",            "Process sale → award points → Clubby benefit"),
    ("GET",  "/api/members/:id/history", "Purchase history (last 10)"),
    ("GET",  "/api/dashboard",           "Revenue, top members, today's stats"),
    ("POST", "/clubby/webhook",          "← Inbound from Clubby (QR scan)"),
    ("POST", "/functions/v1/issue-benefit","← Push benefit to wallet from backend"),
]
for j, (method, path, desc) in enumerate(endpoints):
    by = Inches(2.3 + j * 0.72)
    mc = GREEN if method == "GET" else AMBER
    pill(sl, method, mc, TX, by, pw=Inches(0.65), ph=Inches(0.24))
    t(sl, path, TX+Inches(0.75), by-Inches(0.02), Inches(2.8), Inches(0.3),
      size=11, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    t(sl, desc, TX+Inches(0.75), by+Inches(0.26), Inches(5.3), Inches(0.28),
      size=10, color=GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — SECTION DIVIDER: Platform Analytics
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, TEAL)
rect(sl, 0, 0, Inches(0.07), H, WHITE)
t(sl, "04", Inches(0.3), Inches(0.5), Inches(2), Inches(0.7),
  size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
t(sl, "Platform Analytics", Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.5),
  size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
t(sl, "Cross-business insights for the anchor client",
  Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.7),
  size=20, color=RGBColor(0xCC,0xFF,0xF5), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Platform Analytics Full
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, TEAL)
wordmark(sl)
pgnum(sl, 21, TOTAL)

FW = Inches(6.8); FH = Inches(5.9)
browser_frame(sl, f"{D}/18-portal-platform-analytics.png", Inches(0.3), Inches(0.75), FW, FH)

TX = Inches(7.5)
pill(sl, "ANCHOR CLIENT", TEAL, TX, Inches(0.65))
t(sl, "Cross-business view",
  TX, Inches(1.15), Inches(5.8), Inches(0.7),
  size=28, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.0), Inches(4.5), Inches(0.04), TEAL)
t(sl, "A dedicated 'Platform View' in the portal gives retail chains and anchor clients a live dashboard across all their businesses.",
  TX, Inches(2.15), Inches(5.7), Inches(1.0),
  size=13, color=GRAY, align=PP_ALIGN.LEFT)
stats_demo = [
    ("4", "Businesses on platform"),
    ("12", "Total members"),
    ("13", "Benefits issued"),
    ("62%", "Redemption rate"),
]
for j, (val, label) in enumerate(stats_demo):
    bx = TX + (j % 2) * Inches(2.9)
    by = Inches(3.4 + (j // 2) * 0.85)
    box = sl.shapes.add_shape(1, bx, by, Inches(2.7), Inches(0.72))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = RGBColor(0xE2,0xE8,0xF0); box.line.width = Pt(1)
    t(sl, val, bx+Inches(0.12), by+Inches(0.04), Inches(0.9), Inches(0.38),
      size=22, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
    t(sl, label, bx+Inches(1.05), by+Inches(0.18), Inches(1.55), Inches(0.38),
      size=10, color=GRAY, align=PP_ALIGN.LEFT)

items = [
    "📊  30-day new users + new memberships charts",
    "🏆  Top businesses by member count table",
    "📈  Redemption rates per business",
    "🔄  Live data — refreshes on every page load",
]
for j, b in enumerate(items):
    t(sl, b, TX, Inches(5.3+j*0.44), Inches(5.7), Inches(0.4),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Scenario: 5 customers join + feed
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, GREEN)
wordmark(sl)
pgnum(sl, 22, TOTAL)

# 3 phones side by side
PH = Inches(5.4); PW = Inches(2.55)
phone_frame(sl, f"{S}/P2-dana-wallet.png",  Inches(0.25), Inches(1.0), PW, PH)
phone_frame(sl, f"{S}/P2-yossi-wallet.png", Inches(3.1),  Inches(1.0), PW, PH)
phone_frame(sl, f"{S}/P6-dana-feed.png",    Inches(5.95), Inches(1.0), PW, PH)

TX = Inches(9.0)
pill(sl, "REAL SCENARIO", GREEN, TX, Inches(0.65))
t(sl, "5 customers,\none coffee shop",
  TX, Inches(1.15), Inches(4.1), Inches(1.0),
  size=26, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, TX, Inches(2.25), Inches(3.5), Inches(0.04), GREEN)
t(sl, "End-to-end scenario tested live: owner configures the shop, 5 customers enroll via QR, receive welcome benefits, see the feed, and the owner watches the dashboard update in real time.",
  TX, Inches(2.4), Inches(4.1), Inches(1.5),
  size=12, color=GRAY, align=PP_ALIGN.LEFT)
milestones = [
    "✓  QR scan → benefit in wallet < 2s",
    "✓  Bronze/Silver/Gold tier badges",
    "✓  Feed post from portal → wallet live",
    "✓  Owner notify modal: 'Send to 7 members'",
    "✓  Platform analytics updated in real time",
]
for j, m in enumerate(milestones):
    t(sl, m, TX, Inches(4.1+j*0.52), Inches(4.1), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, GRAY)
wordmark(sl)
pgnum(sl, 23, TOTAL)

t(sl, "Tech Stack", Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.7),
  size=32, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, Inches(0.4), Inches(1.85), Inches(4.0), Inches(0.04), GRAY)

stacks = [
    (GREEN, "Mobile App",
     "Expo SDK 56 + React Native 0.85\nExpo Router v4 (file-based navigation)\nNativeWind v4 + Urbanist font\nTanStack Query (data fetching)\nSupabase JS client"),
    (INDIGO, "Business Portal",
     "React 18 + Vite\nRecharts (analytics charts)\nSupabase JS client\nCSS variables (language/theme)\nFull RTL + i18next"),
    (AMBER, "Backend",
     "Supabase (PostgreSQL 17)\nEdge Functions (Deno)\nRow Level Security (RLS)\nRealtime subscriptions\npg_cron (scheduled jobs)"),
    (TEAL, "Integration",
     "Node.js / Express (mock POS)\nHMAC-SHA256 webhook signing\nExpo Push + FCM / APNs\nWeb Push (VAPID)\nSupabase Storage (images)"),
]
for i, (accent, title, body) in enumerate(stacks):
    cx = Inches(0.3 + i * 3.25)
    card = sl.shapes.add_shape(1, cx, Inches(2.2), Inches(3.1), Inches(4.8))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2,0xE8,0xF0); card.line.width = Pt(1)
    rect(sl, cx, Inches(2.2), Inches(3.1), Inches(0.06), accent)
    t(sl, title, cx+Inches(0.15), Inches(2.38), Inches(2.8), Inches(0.4),
      size=14, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    t(sl, body, cx+Inches(0.15), Inches(2.9), Inches(2.85), Inches(4.0),
      size=12, color=GRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — What's Ready vs Roadmap
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, PINK)
wordmark(sl)
pgnum(sl, 24, TOTAL)

t(sl, "What's built  ·  What's next",
  Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, Inches(0.4), Inches(1.85), Inches(4.5), Inches(0.04), PINK)

built = [
    "✅ Phone OTP auth (customer + portal)",
    "✅ QR scan enrollment + welcome benefit",
    "✅ Wallet: balance, stamps, points, tiers",
    "✅ Discover: GPS, categories, search",
    "✅ Business portal: full CRM + analytics",
    "✅ Promotions, stamp cards, tiers, points",
    "✅ Stories + feed posts",
    "✅ Segmented push notifications",
    "✅ Webhook integration (HMAC signed)",
    "✅ issue-benefit REST edge function",
    "✅ Osher-Ad mock POS + demo flow",
    "✅ Platform (anchor client) analytics",
    "✅ Hebrew RTL + i18n (5 languages)",
    "✅ Dark mode + Urbanist design system",
]
next_up = [
    "🔜 Live SMS OTP (Twilio/Vonage)",
    "🔜 GPS map view in Discover",
    "🔜 Camera barcode scan in manual coupon",
    "🔜 Anchor client multi-chain dashboard",
    "🔜 Digital receipts (POS push)",
    "🔜 AI-driven campaign recommendations",
    "🔜 App Store + Play Store submission",
    "🔜 Production deploy (Vercel + Supabase)",
]

# Built column
rect(sl, Inches(0.4), Inches(2.2), Inches(6.2), Inches(4.8), RGBColor(0xF0,0xFD,0xF4))
t(sl, "Built & working", Inches(0.6), Inches(2.3), Inches(5.8), Inches(0.4),
  size=13, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
for j, b in enumerate(built):
    t(sl, b, Inches(0.6), Inches(2.85+j*0.31), Inches(5.9), Inches(0.28),
      size=11, color=DARK, align=PP_ALIGN.LEFT)

# Next column
rect(sl, Inches(6.9), Inches(2.2), Inches(6.1), Inches(4.8), RGBColor(0xFF,0xF7,0xED))
t(sl, "Coming next", Inches(7.1), Inches(2.3), Inches(5.7), Inches(0.4),
  size=13, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
for j, b in enumerate(next_up):
    t(sl, b, Inches(7.1), Inches(2.85+j*0.52), Inches(5.7), Inches(0.45),
      size=12, color=DARK, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Deployment Plan
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, OFF)
rect(sl, 0, 0, Inches(0.06), H, GREEN)
wordmark(sl)
pgnum(sl, 25, TOTAL)

t(sl, "Production deployment — ~$12/year",
  Inches(0.4), Inches(1.0), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=DARK, align=PP_ALIGN.LEFT)
rect(sl, Inches(0.4), Inches(1.85), Inches(5.5), Inches(0.04), GREEN)
t(sl, "Everything needed to go live is already coded. The deployment plan is fully written.",
  Inches(0.4), Inches(2.05), Inches(12.5), Inches(0.55),
  size=14, color=GRAY, align=PP_ALIGN.LEFT)

steps_deploy = [
    (GREEN,  "1", "Supabase Cloud",    "Create project → db push → deploy edge functions\nEnable pg_cron → set Twilio credentials"),
    (INDIGO, "2", "Vercel (Portal)",   "cd apps/portal && npm run build → vercel deploy\nCustom domain: portal.clubby.com"),
    (AMBER,  "3", "Vercel (App)",      "npx expo export --platform web → vercel deploy\nCustom domain: app.clubby.com"),
    (TEAL,   "4", "Native Apps",       "eas build --platform all --profile production\neas submit → App Store + Play Store"),
]
for i, (accent, num, title, body) in enumerate(steps_deploy):
    cx = Inches(0.3 + i * 3.25)
    card = sl.shapes.add_shape(1, cx, Inches(2.85), Inches(3.1), Inches(4.3))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE2,0xE8,0xF0); card.line.width = Pt(1)
    circle = sl.shapes.add_shape(9, cx+Inches(0.15), Inches(3.0), Inches(0.45), Inches(0.45))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent; circle.line.fill.background()
    t(sl, num, cx+Inches(0.22), Inches(3.04), Inches(0.3), Inches(0.35),
      size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    t(sl, title, cx+Inches(0.7), Inches(3.0), Inches(2.3), Inches(0.45),
      size=13, bold=True, color=DARK, align=PP_ALIGN.LEFT)
    t(sl, body, cx+Inches(0.15), Inches(3.55), Inches(2.85), Inches(1.5),
      size=11, color=GRAY, align=PP_ALIGN.LEFT)

t(sl, "Domain: clubby.com (~$12/yr Cloudflare)\nHTTPS auto-provisioned by Vercel — unlocks camera QR and web push notifications",
  Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
  size=12, color=LGRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill(sl, NAVY)
rect(sl, 0, 0, Inches(0.07), H, GBRT)
wordmark(sl)

t(sl, "Ready to connect\nyour loyalty program.", Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.4),
  size=46, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
rect(sl, W-Inches(4.5), Inches(4.0), Inches(4.2), Inches(0.06), GBRT)
t(sl, "Full stack built. Backend integration live. Demo ready.",
  Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.6),
  size=18, color=LGRAY, align=PP_ALIGN.RIGHT)

# CTA boxes
for i, (label, val) in enumerate([
    ("Customer App", "http://localhost:8081"),
    ("Business Portal", "http://localhost:5173"),
    ("Osher-Ad POS Demo", "http://localhost:3001"),
]):
    bx = Inches(1.5 + i * 3.6)
    box = sl.shapes.add_shape(1, bx, Inches(5.1), Inches(3.4), Inches(1.8))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x14,0x20,0x33)
    box.line.fill.background()
    t(sl, label, bx+Inches(0.15), Inches(5.2), Inches(3.1), Inches(0.35),
      size=12, bold=True, color=GBRT, align=PP_ALIGN.LEFT)
    t(sl, val, bx+Inches(0.15), Inches(5.6), Inches(3.1), Inches(0.35),
      size=10, color=LGRAY, align=PP_ALIGN.LEFT)

prs.save(OUT)
print(f"\n✅ Deck saved: {OUT}")
print(f"   Slides: {len(prs.slides)}")
