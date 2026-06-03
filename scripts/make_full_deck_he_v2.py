"""
Clubby — Full Hebrew Deck v2 (Clubby_Full_HE_v2.pptx)
Complete redesign:
- Large screenshots (55% of slide height) inside device frames
- Bold hero cover with mockup
- Split layout: big visual left, punchy text right
- Dark section dividers
- One claim per slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

S   = "/Users/I560647/Library/CloudStorage/OneDrive-SAPSE/Desktop/finance/Clubby/screenshots"
A   = f"{S}/after-ui-improvements"
D   = f"{S}/demo"   # real-world demo screenshots with Hebrew + real businesses
OUT = "/Users/I560647/Library/CloudStorage/OneDrive-SAPSE/Desktop/finance/Clubby/slides"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x09, 0x0E, 0x1A)   # near-black background
GREEN  = RGBColor(0x2E, 0xCC, 0x71)   # brand green
INDIGO = RGBColor(0x5B, 0x4C, 0xF5)   # section 2 accent
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # section 3 accent
PINK   = RGBColor(0xEC, 0x48, 0x99)   # section 4 accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
OFF    = RGBColor(0xF8, 0xFA, 0xFC)   # slide background (near-white)
GRAY   = RGBColor(0x64, 0x74, 0x8B)
LGRAY  = RGBColor(0x94, 0xA3, 0xB8)
DARK   = RGBColor(0x1E, 0x29, 0x3B)   # dark text

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def set_rtl(p_xml):
    pPr = p_xml.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p_xml, qn('a:pPr'))
        p_xml.insert(0, pPr)
    pPr.set('rtl', '1')

def fill(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def t(slide, text, x, y, w, h, size=18, bold=False, color=DARK,
      align=PP_ALIGN.RIGHT, italic=False, rtl=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    if rtl: set_rtl(p._p)

def rect(slide, x, y, w, h, color, radius=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def img(slide, path, x, y, w, h=None):
    if not path or not os.path.exists(path): return False
    if h: slide.shapes.add_picture(path, x, y, w, h)
    else: slide.shapes.add_picture(path, x, y, w)
    return True

def phone_frame(slide, screenshot_path, x, y, frame_w, frame_h):
    """Place screenshot inside a rounded phone-like frame"""
    if not screenshot_path or not os.path.exists(screenshot_path): return
    # dark frame background
    frm = slide.shapes.add_shape(
        9,   # freeform → we use rounded rect type 9 = RoundRect
        x, y, frame_w, frame_h
    )
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    frm.line.color.rgb = RGBColor(0x33, 0x3A, 0x55)
    frm.line.width = Pt(1.5)
    # screenshot slightly inset
    pad = Inches(0.12)
    slide.shapes.add_picture(screenshot_path, x+pad, y+pad, frame_w-pad*2, frame_h-pad*2)

def web_frame(slide, screenshot_path, x, y, fw, fh):
    """Browser-chrome style frame for portal screenshots"""
    if not screenshot_path or not os.path.exists(screenshot_path): return
    chrome_h = Inches(0.25)
    frm = slide.shapes.add_shape(1, x, y, fw, fh)
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
    frm.line.color.rgb = RGBColor(0xCB, 0xD5, 0xE1); frm.line.width = Pt(1)
    # chrome bar
    chrome = slide.shapes.add_shape(1, x, y, fw, chrome_h)
    chrome.fill.solid(); chrome.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    chrome.line.fill.background()
    # screenshot
    slide.shapes.add_picture(screenshot_path, x, y+chrome_h, fw, fh-chrome_h)

def accent_dot(slide, accent):
    """Small colored accent circle top-right"""
    r = slide.shapes.add_shape(9, W-Inches(0.55), Inches(0.2), Inches(0.18), Inches(0.18))
    r.fill.solid(); r.fill.fore_color.rgb = accent; r.line.fill.background()

def pgnum(slide, n, total, color=LGRAY):
    t(slide, f"{n} / {total}", Inches(0.25), H-Inches(0.4), Inches(1.2), Inches(0.35),
      size=10, color=color, align=PP_ALIGN.LEFT, rtl=False)

def tag_pill(slide, label, accent, x, y):
    """Colored tag pill e.g. 'לקוח' """
    # light background derived from accent
    r = min(255, int(accent[0]*0.15 + 0xFF*0.85))
    g = min(255, int(accent[1]*0.15 + 0xFF*0.85))
    b = min(255, int(accent[2]*0.15 + 0xFF*0.85))
    pill = slide.shapes.add_shape(9, x, y, Inches(0.9), Inches(0.28))
    pill.fill.solid(); pill.fill.fore_color.rgb = RGBColor(r, g, b)
    pill.line.color.rgb = accent; pill.line.width = Pt(0.75)
    t(slide, label, x+Inches(0.05), y+Inches(0.02), Inches(0.8), Inches(0.26),
      size=9, bold=True, color=accent, align=PP_ALIGN.CENTER, rtl=True)

WORDMARK = "/Users/I560647/Library/CloudStorage/OneDrive-SAPSE/Desktop/finance/Clubby/logo/wordmark.png"

# ── Slide builders ────────────────────────────────────────────────────────────

def make_cover(slide, headline, subline, bullets, accent, logo_text="🏪"):
    """Dark hero cover with large text and visual band"""
    fill(slide, NAVY)
    # Left accent stripe
    rect(slide, 0, 0, Inches(0.06), H, accent)
    # Wordmark logo top-left
    if os.path.exists(WORDMARK):
        slide.shapes.add_picture(WORDMARK, Inches(0.25), Inches(0.18), Inches(2.4), Inches(0.65))
    else:
        t(slide, "clubby.", Inches(0.3), Inches(0.25), Inches(2.5), Inches(0.5),
          size=22, bold=True, color=GREEN, align=PP_ALIGN.LEFT, rtl=False)
    # Big headline — right aligned (RTL)
    t(slide, headline, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.0),
      size=52, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    # Accent line
    rect(slide, W-Inches(3.5), Inches(3.8), Inches(3.2), Inches(0.05), accent)
    # Subline
    t(slide, subline, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.7),
      size=20, color=LGRAY, align=PP_ALIGN.RIGHT)
    # Bullets in a bottom band
    rect(slide, 0, Inches(5.0), W, Inches(2.0), RGBColor(0x0F,0x17,0x2A))
    for j, b in enumerate(bullets):
        col = j * 3.3
        t(slide, b, Inches(0.3 + col), Inches(5.2), Inches(3.0), Inches(0.55),
          size=14, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

def make_section(slide, label, accent, sub=""):
    fill(slide, accent)
    rect(slide, 0, 0, Inches(0.06), H, WHITE)
    t(slide, label, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.8),
      size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        t(slide, sub, Inches(2.0), Inches(4.5), Inches(9.3), Inches(0.8),
          size=18, color=RGBColor(0xFF,0xFF,0xFF,), align=PP_ALIGN.CENTER)

def make_mobile_slide(slide, num, total, title, subtitle, bullets, img_path, accent, tag=""):
    """Phone screenshot left (big), text right"""
    fill(slide, OFF)
    # Left accent stripe
    rect(slide, 0, 0, Inches(0.05), H, accent)

    # Phone frame — occupies most of left side
    PH = Inches(6.5); PW = Inches(3.3)
    PX = Inches(0.45); PY = (H - PH) / 2
    phone_frame(slide, img_path, PX, PY, PW, PH)

    # Text area
    TX = Inches(4.2); TW = Inches(8.8)
    # Tag pill
    if tag:
        tag_pill(slide, tag, accent, TX, Inches(0.65))
    # Title
    t(slide, title, TX, Inches(1.1), TW-Inches(0.3), Inches(1.0),
      size=32, bold=True, color=DARK)
    # Accent line
    rect(slide, W-Inches(3.5), Inches(2.25), Inches(3.2), Inches(0.04), accent)
    # Subtitle
    t(slide, subtitle, TX, Inches(2.4), TW-Inches(0.3), Inches(1.4),
      size=16, color=GRAY)
    # Bullets
    for j, b in enumerate(bullets):
        # bullet dot
        dot = slide.shapes.add_shape(9, W-Inches(0.55), Inches(4.05+j*0.65)-Inches(0.02),
                                      Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        t(slide, b, TX, Inches(3.9 + j*0.65), TW-Inches(0.75), Inches(0.55),
          size=15, color=DARK)
    pgnum(slide, num, total, GRAY)

def make_portal_slide(slide, num, total, title, subtitle, bullets, img_path, accent, tag=""):
    """Web screenshot left (wide), text right"""
    fill(slide, OFF)
    rect(slide, 0, 0, Inches(0.05), H, accent)

    # Web frame — wider
    WFW = Inches(5.8); WFH = Inches(5.6)
    WFX = Inches(0.35); WFY = Inches(0.85)
    web_frame(slide, img_path, WFX, WFY, WFW, WFH)

    TX = Inches(6.5); TW = Inches(6.6)
    if tag:
        tag_pill(slide, tag, accent, TX, Inches(0.65))
    t(slide, title, TX, Inches(1.1), TW-Inches(0.2), Inches(1.0),
      size=30, bold=True, color=DARK)
    rect(slide, W-Inches(3.2), Inches(2.25), Inches(2.9), Inches(0.04), accent)
    t(slide, subtitle, TX, Inches(2.4), TW-Inches(0.2), Inches(1.4),
      size=15, color=GRAY)
    for j, b in enumerate(bullets):
        dot = slide.shapes.add_shape(9, W-Inches(0.5), Inches(4.0+j*0.65)-Inches(0.02),
                                      Inches(0.12), Inches(0.12))
        dot.fill.solid(); dot.fill.fore_color.rgb = accent; dot.line.fill.background()
        t(slide, b, TX, Inches(3.85 + j*0.65), TW-Inches(0.65), Inches(0.55),
          size=14, color=DARK)
    pgnum(slide, num, total, GRAY)

def make_text_slide(slide, num, total, title, subtitle, bullets, accent, tag=""):
    """No image — large text, stat callouts"""
    fill(slide, OFF)
    rect(slide, 0, 0, Inches(0.05), H, accent)
    if tag:
        tag_pill(slide, tag, accent, Inches(0.4), Inches(0.65))
    t(slide, title, Inches(0.4), Inches(1.1), Inches(12.5), Inches(1.0),
      size=36, bold=True, color=DARK)
    rect(slide, W-Inches(4.5), Inches(2.25), Inches(4.2), Inches(0.04), accent)
    t(slide, subtitle, Inches(0.4), Inches(2.4), Inches(12.5), Inches(1.2),
      size=18, color=GRAY)
    cols = 2
    for j, b in enumerate(bullets):
        row, col = divmod(j, cols)
        bx = W - Inches(0.5) - (col+1)*Inches(6.2) + Inches(0.2)
        by = Inches(3.9) + row * Inches(0.88)
        # Card-style bullet
        card = slide.shapes.add_shape(1, bx, by, Inches(5.8), Inches(0.72))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xE2,0xE8,0xF0); card.line.width = Pt(0.75)
        t(slide, b, bx+Inches(0.15), by+Inches(0.12), Inches(5.5), Inches(0.5),
          size=15, color=DARK)
    pgnum(slide, num, total, GRAY)

def make_closing(slide, headline, sub, bullets, accent):
    fill(slide, NAVY)
    rect(slide, 0, 0, Inches(0.06), H, accent)
    if os.path.exists(WORDMARK):
        slide.shapes.add_picture(WORDMARK, Inches(0.25), Inches(0.18), Inches(2.4), Inches(0.65))
    else:
        t(slide, "clubby.", Inches(0.3), Inches(0.25), Inches(2.5), Inches(0.5),
          size=22, bold=True, color=GREEN, align=PP_ALIGN.LEFT, rtl=False)
    t(slide, headline, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.0),
      size=40, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    rect(slide, W-Inches(3.5), Inches(3.6), Inches(3.2), Inches(0.05), accent)
    t(slide, sub, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.8),
      size=18, color=LGRAY, align=PP_ALIGN.RIGHT)
    band = slide.shapes.add_shape(1, 0, Inches(4.9), W, Inches(2.1), )
    band.fill.solid(); band.fill.fore_color.rgb = RGBColor(0x0F,0x17,0x2A)
    band.line.fill.background()
    for j, b in enumerate(bullets):
        t(slide, b, Inches(0.5 + j*4.4), Inches(5.1), Inches(4.0), Inches(0.75),
          size=15, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.RIGHT)

# ── SLIDE DATA ─────────────────────────────────────────────────────────────────
content_items = []  # (type, *args)

def C(typ, *args): content_items.append((typ, args))

# Cover
C("cover", "ארנק הנאמנות שמחבר לקוחות ועסקים",
  "הצטרפות בסריקת QR אחת | ארנק דיגיטלי | ניתוח נתונים לרשתות ענק",
  ["📱 ארנק לקוח דיגיטלי", "🏪 פורטל ניהול לעסק", "📊 אנליטיקה לרשת קמעונאות"],
  GREEN)

# Customer section
C("section", "🛍️ חוויית הלקוח", GREEN, "מהסריקה הראשונה ועד לפרס — חלק, מהיר, ומהנה")

C("mobile", "1. גישה מהירה — ללא הורדה",
  "הלקוח נכנס לאפליקציה דרך הדפדפן בטלפון או מוריד אפליקציה. הרשמה תוך 30 שניות עם מספר טלפון בלבד.",
  ["🌐 app.clubby.com — ללא הורדה מחנות האפליקציות",
   "📲 אפליקציה ל-iOS ו-Android",
   "📞 הרשמה עם מספר טלפון — ללא סיסמה",
   "⚡ מגישה לארנק פעיל תוך 30 שניות"],
  f"{D}/08-signin-he.png", GREEN, "לקוח")

C("mobile", "2. הצטרף בסריקת QR אחת",
  "קוד QR מוצג בחנות. הלקוח מכוון את מצלמת הטלפון — ומצטרף מיידית. ללא חיכוך, ללא טפסים.",
  ["📷 עובד עם מצלמת הטלפון הרגילה — ללא אפליקציית סריקה",
   "🎁 הטבת ברוכים הבאים — אם בעל העסק הגדיר מבצע קבלת פנים",
   "🏢 כל סניף יכול לקבל קוד QR נפרד",
   "🔄 סריקה חוזרת בטוחה — ללא כפילויות"],
  f"{D}/02-pizza-store-he.png", GREEN, "לקוח")

C("mobile", "3. ארנק — הכל במקום אחד",
  "המועדונים, החותמות, הנקודות, הדרגות, פיד המבצעים וההטבות הפעילות — ממשק יוקרתי ומהיר.",
  ["💰 יתרת קרדיט + היסטוריה בלחיצה",
   "📣 פיד מבצעים בסגנון WhatsApp",
   "🥇 כרטיסיות חותמות + יתרת נקודות",
   "🥈 תג דרגה על כרטיס המועדון"],
  f"{D}/01-wallet-he.png", GREEN, "לקוח")

C("mobile", "4. פיד מבצעים — בסגנון WhatsApp",
  "עסקים שולחים הכרזות ומבצעים שמופיעים ישירות בארנק הלקוח כפיד גולל — בזמן אמת.",
  ["📣 4 סוגי פוסטים: מבצע, הצעה, הכרזה, סטורי",
   "🎨 תווית צבעונית לכל סוג",
   "🔘 כפתור CTA מפנה לפרופיל החנות",
   "⏰ פוסטים פגים אוטומטית בתאריך שנקבע"],
  f"{D}/01-wallet-he.png", GREEN, "לקוח")

C("mobile", "5. גלה עסקים — לפי מיקום",
  "עסקים ממוינים לפי מרחק GPS, סינון קטגוריה, 'פתוח עכשיו' וחיפוש חופשי.",
  ["📍 'קרוב אלי' — ממוין לפי מרחק",
   "🟢 סינון 'פתוח עכשיו'",
   "🔍 חיפוש לפי שם או קטגוריה",
   "➕ הצטרף מיידית בלחיצה אחת"],
  f"{D}/04-discover-he.png", GREEN, "לקוח")

C("mobile", "6. חותמות + נקודות + דרגות",
  "שלושה מנגנוני נאמנות מקבילים — כל אחד מוסיף שכבה של מוטיבציה לחזור.",
  ["🟢 כרטיסיית חותמות עם התקדמות בזמן אמת",
   "⭐ נקודות שנצברות בכל סריקת QR",
   "🥈 תג דרגה (ברונזה/כסף/זהב) על הארנק",
   "🎁 פרס אוטומטי בהשלמת הכרטיסייה"],
  f"{D}/03-supermarket-store-he.png", GREEN, "לקוח")

C("mobile", "7. מימוש הטבה — בלחיצה אחת",
  "לוחץ 'השתמש', מראה למוכר ולוחץ אשר. מסונכרן בזמן אמת עם פורטל העסק.",
  ["👁 הצג למוכר ולחץ אשר",
   "✅ מסומן כמומש בשני הצדדים מיידית",
   "🔄 סנכרון בזמן אמת — WebSocket",
   "📋 מופיע בהיסטוריית ההטבות"],
  f"{D}/06-notifications-he.png", GREEN, "לקוח")

C("mobile", "8. מצב כהה | עברית מלאה | RTL",
  "ממשק מלא בעברית עם RTL, מצב כהה אוטומטי, ותמיכה ב-iOS, Android ו-Web.",
  ["🌙 מצב כהה: אוטומטי / בהיר / כהה",
   "🌐 עברית מלאה עם פריסת RTL",
   "💾 ההעדפות נשמרות בין פתיחות",
   "📲 iOS, Android, Web — אותה חוויה"],
  f"{D}/07-profile-he.png", GREEN, "לקוח")

# Owner section
C("section", "🏪 פורטל בעל העסק", INDIGO, "מלוח הבקרה ועד לפיד — הכל בלחיצה אחת")

C("portal", "9. הקמה תוך 5 דקות",
  "נרשמים עם אימייל, יוצרים פרופיל עסק, וקוד QR מוכן לסריקה — הכל בפחות מ-5 דקות.",
  ["⚡ הרשמה → קוד QR מוכן ב-5 דקות",
   "🖼 לוגו + תמונת כיסוי + שעות פתיחה",
   "🏢 ריבוי סניפים — קוד QR נפרד לכל סניף",
   "🌙 מצב כהה + עברית בפורטל"],
  f"{D}/p06-qr-he.png", INDIGO, "עסק")

C("portal", "10. ניהול מבצעים | חותמות | נקודות | דרגות",
  "כל מנגנוני הנאמנות במקום אחד — כולל כרטיסיות חותמות, תוכנית נקודות ודרגות.",
  ["🎁 מבצעים: הנחה / קרדיט / פריט חינמי",
   "🥇 כרטיסיות חותמות + הוספת חותמת מהפורטל",
   "⭐ תוכנית נקודות עם פרסים מוגדרים",
   "🏅 דרגות — סף מוגדר על ידי בעל העסק"],
  f"{D}/p03-stamps-he.png", INDIGO, "עסק")

C("portal", "11. פיד פוסטים — WhatsApp לעסק",
  "יצירה, עריכה ומחיקה של פוסטים שיופיעו בארנק הלקוח. הכל בממשק פשוט.",
  ["📣 4 סוגי פוסטים עם צבע ייעודי",
   "🖼 תמונה + כפתור CTA + תאריך תפוגה",
   "✏️ עריכה ומחיקה בכל עת",
   "🔄 פוסט חדש מופיע בארנק הלקוח מיידית"],
  f"{D}/p04-feed-he.png", INDIGO, "עסק")

C("portal", "12. חברים + מימוש מהפורטל",
  "רשימת חברים מלאה, לחיצה לפרטים, היסטוריית הטבות — ואישור מימוש ישירות מהפורטל.",
  ["👥 רשימה: שם, טלפון, תאריך הצטרפות",
   "🔍 לחיצה → היסטוריית הטבות מלאה",
   "✓ מימוש הטבה ישירות מהפורטל",
   "🔄 סנכרון בזמן אמת עם ארנק הלקוח"],
  f"{D}/p02-members-he.png", INDIGO, "עסק")

C("portal", "13. התראות ממוקדות + ייצוא CSV",
  "שלח Push לקהל הנכון: חדשים / פעילים / לא פעילים. מנגנון מעורבות מחדש אוטומטי.",
  ["📣 4 קהלים: הכל / חדשים / פעילים / לא פעילים",
   "🤖 מעורבות מחדש — Push אוטומטי אחרי 30 יום",
   "🎂 פרס יום הולדת אוטומטי",
   "⬇️ ייצוא CSV: חברים + היסטוריית הטבות"],
  f"{D}/p02-members-he.png", INDIGO, "עסק")

C("portal", "14. לוח בקרה + ניתוח נתונים",
  "4 מדדים בזמן אמת, גרפי מגמות ל-30 ימים, וייצוא מלא של כל פעילות ההטבות.",
  ["📊 גידול חברים + מימושים יומיים — 30 יום",
   "💡 טיפ חכם לפי נתוני העסק",
   "⬇️ ייצוא: כל הטבה עם שם, סוג, ערך, מקור",
   "🔄 נתונים בזמן אמת"],
  f"{D}/p01-dashboard-he.png", INDIGO, "עסק")

# Anchor section
C("section", "🏢 לקוח העוגן — פלטפורמת הדאטה", AMBER,
  "נתוני פלטפורמה בזמן אמת לרשת קמעונאות")

C("portal", "15. לוח בקרה של לקוח העוגן",
  "הרשת הגדולה מקבלת גישה לנתונים מצרפיים — כלל המשתמשים, עסקים, הטבות ומגמות — בזמן אמת.",
  ["🏢 כלל העסקים + כלל המשתמשים בפלטפורמה",
   "📊 גרפי גידול משתמשים וחברויות — 30 יום",
   "🏆 10 עסקים מובילים לפי חברים + שיעור מימוש",
   "💡 הדאטה יושב אצלנו — הרשת מקבלת תובנות"],
  f"{D}/p05-anchor-he.png", AMBER, "עוגן")

# Strategy section
C("section", "🚀 אסטרטגיה ו-GTM", PINK, "כיצד מייצרים טראפיק ממאות אלפי משתמשים ביום")

C("text", "16. פתרון הביצה והתרנגולת",
  "בשביל טראפיק צריך עסקים. בשביל עסקים צריך טראפיק. הפתרון: לקוח עוגן שמביא את שניהם.",
  ["🏪 רשת קמעונאות גדולה מקבלת מערכת דאטה מתקדמת — בחינם",
   "📊 הדאטה נמצא בשרתים שלנו — אנחנו שולטים בנרטיב",
   "📱 המשתמשים של הרשת → בסיס הטראפיק של הפלטפורמה",
   "🤝 הרשת מקבלת אחוזים — שיתוף פעולה אסטרטגי עמוק"],
  PINK, "אסטרטגיה")

C("text", "17. זרימת משתמש — אפס חיכוך",
  "הלקוח סורק QR בחנות → מצטרף → קבלות והטבות נוחתות בארנק. ללא הורדת אפליקציה.",
  ["1️⃣ סריקת QR עם מצלמה רגילה — לא צריך אפליקציית סריקה",
   "2️⃣ הרשמה: מספר טלפון + OTP — 30 שניות",
   "3️⃣ הטבת קבלת פנים נוחתת בארנק — אם בעל העסק הגדיר",
   "4️⃣ כל ביקור עתידי — חותמות/נקודות אוטומטיים"],
  PINK, "אסטרטגיה")

# Next steps
C("section", "📋 צעדים הבאים", GREEN, "מהדמו להשקה מסחרית")

C("text", "18. מוכן לדמו — מה בונים הבא",
  "הפלטפורמה מוכנה לדמו מלא. אלה הצעדים לפני השקה מסחרית.",
  ["✅ ארנק לקוח + פורטל עסק + לוח בקרה עוגן — מוכנים",
   "🚀 פרסום App Store + Google Play (EAS Build)",
   "⏳ אינטגרציית POS — API לחיבור מערכות קופה",
   "⏳ AI אנליטיקה — מסעות פרסום חכמים מבוססי דאטה"],
  GREEN, "הבא")

C("text", "19. רודמפ — רבעון קרוב",
  "סדר הפעולות לבניית המוצר המסחרי המלא.",
  ["1️⃣ הפקה ל-App Store + Google Play",
   "2️⃣ העלאה ל-Supabase Cloud + Vercel (Production)",
   "3️⃣ Pilot עם לקוח עוגן ראשון — רשת קמעונאות",
   "4️⃣ אינטגרציית POS + קבלות דיגיטליות"],
  GREEN, "רודמפ")

# Closing
C("closing",
  "קלאבי — הנאמנות שמגיעה לכל לקוח",
  "מהסריקה הראשונה ועד לניתוח הדאטה — פלטפורמה אחת, ערך לכולם.",
  ["📱 לקוח — ארנק דיגיטלי חינמי",
   "💼 עסק — פורטל ניהול מלא",
   "🏢 רשת — תובנות דאטה אסטרטגיות"],
  GREEN)

# ─── RENDER ───────────────────────────────────────────────────────────────────
content_slides = [c for c in content_items if c[0] in ("mobile","portal","text")]
total_content = len(content_slides)
content_counter = 0

for item in content_items:
    typ = item[0]; args = item[1]
    slide = prs.slides.add_slide(blank)

    if typ == "cover":
        make_cover(slide, *args)
    elif typ == "section":
        make_section(slide, *args)
    elif typ == "mobile":
        content_counter += 1
        title, subtitle, bullets, img_path, accent, tag = args
        make_mobile_slide(slide, content_counter, total_content,
                          title, subtitle, bullets, img_path, accent, tag)
    elif typ == "portal":
        content_counter += 1
        title, subtitle, bullets, img_path, accent, tag = args
        make_portal_slide(slide, content_counter, total_content,
                          title, subtitle, bullets, img_path, accent, tag)
    elif typ == "text":
        content_counter += 1
        title, subtitle, bullets, accent, tag = args
        make_text_slide(slide, content_counter, total_content,
                        title, subtitle, bullets, accent, tag)
    elif typ == "closing":
        make_closing(slide, *args)

out = f"{OUT}/Clubby_Full_HE_v2.pptx"
prs.save(out)
print(f"✅ Saved: {out}  ({len(content_items)} slides, {total_content} content slides)")
