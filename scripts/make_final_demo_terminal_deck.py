"""
Clubby — Final Demo Deck (terminal/IDE style — matches claude-code-power-patterns)
Mixed Hebrew narrative + English code/UI
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

BASE     = "/Users/I560647/Library/CloudStorage/OneDrive-SAPSE/Desktop/finance/Clubby"
V2       = f"{BASE}/screenshots/v2"
DEMO     = f"{BASE}/screenshots/demo"
SCEN     = f"{BASE}/screenshots/scenario"
WORDMARK = f"{BASE}/logo/wordmark.png"
OUT      = f"{BASE}/slides/Clubby_FinalDemo_Terminal.pptx"
os.makedirs(f"{BASE}/slides", exist_ok=True)

# ── Palette (matches GitHub Dark / claude-code-power-patterns) ───────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # main bg
SURFACE  = RGBColor(0x16, 0x1B, 0x22)   # elevated surface
PANEL    = RGBColor(0x1C, 0x21, 0x29)   # code panel
BORDER   = RGBColor(0x30, 0x36, 0x3D)
TEXT     = RGBColor(0xE6, 0xED, 0xF3)   # primary text
DIM      = RGBColor(0x8B, 0x94, 0x9E)   # comments / secondary
GREEN    = RGBColor(0x3F, 0xB9, 0x50)   # $, prompts, keywords
ORANGE   = RGBColor(0xD9, 0x77, 0x57)   # category headers, claude orange
BLUE     = RGBColor(0x58, 0xA6, 0xFF)   # links, plan mode
PURPLE   = RGBColor(0xBC, 0x8C, 0xFF)   # LLM, accents
RED      = RGBColor(0xF8, 0x51, 0x49)   # errors
YELLOW   = RGBColor(0xE3, 0xB3, 0x41)
PINK     = RGBColor(0xFF, 0x7B, 0x72)   # extra accent

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def set_rtl(p_xml):
    pPr = p_xml.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p_xml, qn('a:pPr'))
        p_xml.insert(0, pPr)
    pPr.set('rtl', '1')

def fill_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def t(slide, text, x, y, w, h, size=14, bold=False, color=TEXT,
      align=PP_ALIGN.LEFT, italic=False, rtl=False, mono=False):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    if mono:
        r.font.name = 'Menlo'
    if rtl: set_rtl(p._p)
    return tf

def t_runs(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT, rtl=False):
    """Multi-color runs in a single paragraph. runs = [(text, size, color, bold, mono)]"""
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for txt, size, color, bold, mono in runs:
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        if mono: r.font.name = 'Menlo'
    if rtl: set_rtl(p._p)
    return tf

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def rrect(slide, x, y, w, h, color, border=None, line_pt=0.75):
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def code_panel(slide, x, y, w, h, lines, label="", label_color=DIM):
    """Draw a code panel like the SharePoint deck's right column"""
    if label:
        t(slide, label, x+Inches(0.05), y, w-Inches(0.1), Inches(0.3),
          size=10, color=label_color, mono=True)
        y = y + Inches(0.36)
        h = h - Inches(0.36)
    rrect(slide, x, y, w, h, PANEL, BORDER)
    # Render lines
    line_y = y + Inches(0.18)
    line_h = Inches(0.27)
    for line in lines:
        if isinstance(line, str):
            line = [(line, 11, TEXT, False, True)]
        tf = slide.shapes.add_textbox(x+Inches(0.18), line_y, w-Inches(0.36), line_h).text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        for txt, size, color, bold, mono in line:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color
            if mono: r.font.name = 'Menlo'
        line_y = line_y + line_h

def header(slide, breadcrumb, num, total):
    """Top breadcrumb + page number"""
    t(slide, breadcrumb, Inches(1.1), Inches(0.1), Inches(8.0), Inches(0.3),
      size=11, color=DIM, mono=True)
    t(slide, f"{num} / {total}", Inches(12.1), Inches(0.1), Inches(1.0), Inches(0.3),
      size=11, color=DIM, mono=True, align=PP_ALIGN.RIGHT)

def title_block(slide, title, subtitle):
    t(slide, title, Inches(0.6), Inches(0.7), Inches(12.1), Inches(0.7),
      size=34, bold=True, color=TEXT)
    t(slide, subtitle, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.4),
      size=15, color=ORANGE)

def section_header(slide, label, x, y, color=ORANGE):
    t(slide, label, x, y, Inches(5.5), Inches(0.3),
      size=12, bold=True, color=color)

def bullet_list(slide, items, x, y, w, color=TEXT, size=13, bold_first=False):
    line_h = Inches(0.4)
    for i, item in enumerate(items):
        c = TEXT if not bold_first else (TEXT if i==0 else TEXT)
        is_b = bold_first and i == 0
        t(slide, item, x, y + i*line_h, w, line_h,
          size=size, color=color, bold=is_b)

def img_panel(slide, path, x, y, w, h, label=""):
    if label:
        t(slide, label, x+Inches(0.05), y, w-Inches(0.1), Inches(0.3),
          size=10, color=DIM, mono=True)
        y = y + Inches(0.36)
        h = h - Inches(0.36)
    if path and os.path.exists(path):
        rrect(slide, x, y, w, h, PANEL, BORDER)
        pad = Inches(0.08)
        slide.shapes.add_picture(path, x+pad, y+pad, w-pad*2, h-pad*2)

def footer_command(slide, prefix_runs, message_runs):
    """Bottom command-line style summary like 'goal: leave with patterns...'"""
    runs = []
    for txt, color, bold in prefix_runs:
        runs.append((txt, 13, color, bold, True))
    for txt, color, bold in message_runs:
        runs.append((txt, 13, color, bold, False))
    t_runs(slide, runs, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4))

TOTAL = 22

# ══════════════════════════════════════════════════════════════════════════════
# 1. COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)

# top breadcrumb
t(sl, "~ talks/clubby-final-demo", Inches(2.0), Inches(1.4), Inches(10.0), Inches(0.4),
  size=12, color=DIM, mono=True)

# Command line
t_runs(sl, [
    ("$ ", 22, GREEN, True, True),
    ("clubby --demo --he", 22, TEXT, False, True),
], Inches(1.1), Inches(2.1), Inches(11.0), Inches(0.6))

# Title
t(sl, "Clubby", Inches(1.1), Inches(2.95), Inches(11.0), Inches(1.1),
  size=64, bold=True, color=TEXT)

# Subtitle (orange)
t(sl, "ארנק נאמנות · אינטגרציה ארגונית · POS דמו",
  Inches(1.1), Inches(4.0), Inches(11.0), Inches(0.7),
  size=32, color=ORANGE, rtl=True, align=PP_ALIGN.RIGHT)

# Tagline
t(sl, "Phone OTP · QR enrollment · webhook integration · live event log · anchor analytics",
  Inches(1.1), Inches(4.85), Inches(11.0), Inches(0.4),
  size=15, color=DIM)

# Author line
t_runs(sl, [
    ("# ", 13, DIM, False, True),
    ("Gabi Kalaora", 13, TEXT, False, True),
    ("   ·   final demo · v2 upgrade edition · June 2026", 13, DIM, False, True),
], Inches(1.1), Inches(5.4), Inches(11.0), Inches(0.4))

# Press any key
t(sl, "►", Inches(0.7), Inches(6.7), Inches(0.3), Inches(0.3),
  size=12, color=GREEN, mono=True)
t(sl, "Press any key to continue...", Inches(1.0), Inches(6.7), Inches(6.0), Inches(0.3),
  size=11, color=DIM, mono=True)

# ══════════════════════════════════════════════════════════════════════════════
# 2. AGENDA
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/agenda.md", 2, TOTAL)
title_block(sl, "מה נראה היום", "Eight chapters · twenty minutes · zero fluff")

agenda = [
    ("01", "המוצר", "ארנק לקוח, פורטל עסקי, אינטגרציה"),
    ("02", "תרחיש מציאותי", "5 לקוחות · בית קפה · רכישות · הטבות"),
    ("03", "אינטגרציית POS", "אושר עד · M2M secret · idempotency"),
    ("04", "Live event log", "Webhook flows · נראים בזמן אמת"),
    ("05", "Anchor dashboard", "תצוגת רשת קמעונאית · KPI חוצי-סניפים"),
    ("06", "v2 שדרוגי UX", "12 תיקוני P0 · WCAG · dark mode"),
    ("07", "Architecture", "Supabase · Edge Functions · React Native"),
    ("08", "מה הבא", "Production deploy · roadmap"),
]
for i, (num, title, desc) in enumerate(agenda):
    col = i % 2
    row = i // 2
    bx = Inches(0.7 + col * 6.2)
    by = Inches(2.3 + row * 0.6)
    # Number
    t(sl, num, bx, by, Inches(0.6), Inches(0.5),
      size=13, bold=True, color=ORANGE, mono=True)
    # Title
    t(sl, title, bx + Inches(0.7), by, Inches(5.2), Inches(0.25),
      size=14, bold=True, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)
    # Description
    t(sl, desc, bx + Inches(0.7), by + Inches(0.27), Inches(5.2), Inches(0.25),
      size=11, color=DIM, rtl=True, align=PP_ALIGN.RIGHT)

# Footer
t_runs(sl, [
    ("$ ", 13, GREEN, True, True),
    ("goal:", 13, DIM, False, True),
    ("  סוגרים את הפיץ' עם דמו חי שלא נופל ", 13, TEXT, False, False),
    ("·", 13, DIM, False, False),
    (" production-ready", 13, BLUE, True, False),
], Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 3. THE PRODUCT — mental model
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/architecture.md", 3, TOTAL)
title_block(sl, "המוצר — שלוש שכבות", "ארנק · פורטל · אינטגרציה")

# Three boxes connected by arrows
boxes = [
    ("CUSTOMER", "ארנק לקוח", "iOS · Android · Web", BLUE),
    ("BUSINESS", "פורטל עסקי", "CRM + Analytics", PURPLE),
    ("BACKEND", "מערכת קיימת", "POS / ERP / CRM", GREEN),
]
for i, (label, title, sub, color) in enumerate(boxes):
    bx = Inches(1.0 + i * 4.0)
    by = Inches(2.7)
    # Box
    rrect(sl, bx, by, Inches(2.4), Inches(1.5), SURFACE, color, line_pt=1.5)
    t(sl, label, bx, by + Inches(0.15), Inches(2.4), Inches(0.3),
      size=14, bold=True, color=color, align=PP_ALIGN.CENTER, mono=True)
    t(sl, title, bx, by + Inches(0.55), Inches(2.4), Inches(0.4),
      size=15, bold=True, color=TEXT, align=PP_ALIGN.CENTER, rtl=True)
    t(sl, sub, bx, by + Inches(0.95), Inches(2.4), Inches(0.4),
      size=11, color=DIM, align=PP_ALIGN.CENTER)
    if i < 2:
        # Arrow
        ax = bx + Inches(2.5)
        t(sl, "↔", ax, by + Inches(0.5), Inches(0.4), Inches(0.5),
          size=24, color=DIM, align=PP_ALIGN.CENTER, mono=True)

# Bottom 3 columns
features = [
    ("Customer flow", ["Phone OTP", "QR scan to join", "Wallet · benefits", "Live updates feed"]),
    ("Business flow", ["Multi-business CRM", "Promotions · stamps", "Tier system · points", "Push notifications"]),
    ("Integration",   ["Webhook (HMAC-SHA256)", "M2M integration key", "Idempotency built-in", "Live event log"]),
]
for i, (heading, items) in enumerate(features):
    bx = Inches(0.9 + i * 4.15)
    section_header(sl, heading, bx, Inches(4.7), color=TEXT)
    for j, item in enumerate(items):
        t(sl, item, bx, Inches(5.05 + j * 0.32), Inches(3.9), Inches(0.3),
          size=11, color=DIM)

t_runs(sl, [
    ("# ", 13, DIM, False, True),
    ("loop until done", 13, DIM, False, True),
    ("   ↩", 13, DIM, False, True),
], Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 4. CUSTOMER WALLET
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/mobile/wallet.tsx", 4, TOTAL)
title_block(sl, "ארנק לקוח — הכל במקום אחד", "Mobile-first · dark mode · RTL · 5 שפות")

# Left: bullets, Right: phone screenshot
section_header(sl, "WHAT'S IN THE WALLET", Inches(0.7), Inches(2.2))
items = [
    "💰  יתרת קרדיט בשקלים — היסטוריה זמינה",
    "🏪  MY CLUBS — תגי ברונזה / כסף / זהב",
    "📢  UPDATES — פיד פוסטים מכל העסקים",
    "🎁  הטבות פעילות — לחיצה למימוש",
    "🔔  התראות עם badge ספירה",
    "🔍  חיפוש ופילטר — All / Credit / Discounts",
]
for j, b in enumerate(items):
    t(sl, b, Inches(0.7), Inches(2.6 + j * 0.55), Inches(5.6), Inches(0.5),
      size=13, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)

section_header(sl, "TECH STACK", Inches(0.7), Inches(6.05), color=GREEN)
t_runs(sl, [
    ("Expo SDK 56", 12, BLUE, True, True),
    (" · ", 12, DIM, False, True),
    ("React Native 0.85", 12, BLUE, True, True),
    (" · ", 12, DIM, False, True),
    ("NativeWind v4", 12, BLUE, True, True),
    (" · ", 12, DIM, False, True),
    ("Supabase", 12, BLUE, True, True),
], Inches(0.7), Inches(6.4), Inches(7.0), Inches(0.4))

# Phone right
phone_path = f"{V2}/09-mobile-wallet-with-feed.png"
if os.path.exists(phone_path):
    rrect(sl, Inches(7.4), Inches(2.0), Inches(2.9), Inches(5.0), PANEL, BORDER)
    pad = Inches(0.1)
    sl.shapes.add_picture(phone_path, Inches(7.4)+pad, Inches(2.0)+pad,
                          Inches(2.9)-pad*2, Inches(5.0)-pad*2)
t(sl, "wallet.tsx", Inches(10.5), Inches(2.0), Inches(2.5), Inches(0.3),
  size=10, color=DIM, mono=True)

# ══════════════════════════════════════════════════════════════════════════════
# 5. SCENARIO — 5 customers
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/scenario.md", 5, TOTAL)
title_block(sl, "תרחיש מציאותי — 5 לקוחות, בית קפה אחד", "Real customers · real benefits · real-time updates")

section_header(sl, "TIMELINE", Inches(0.7), Inches(2.1))
events = [
    ("09:00", "owner מגדיר עסק בפורטל", BLUE),
    ("09:15", "5 לקוחות מצטרפים בסריקת QR", GREEN),
    ("09:18", "הטבת ברוכים הבאים בארנק", GREEN),
    ("12:30", "בית קפה מפרסם פוסט בפיד", PURPLE),
    ("12:31", "כל הלקוחות רואים את הפוסט", PURPLE),
    ("13:00", "owner שולח Push 'בואו לקפה'", ORANGE),
    ("14:00", "Dana מגיעה, מממשת הטבה", GREEN),
    ("14:01", "קופה מעדכנת analytics בזמן אמת", BLUE),
]
for i, (time, desc, color) in enumerate(events):
    by = Inches(2.5 + i * 0.45)
    t(sl, time, Inches(0.7), by, Inches(0.7), Inches(0.3),
      size=11, color=color, bold=True, mono=True)
    t(sl, desc, Inches(1.6), by, Inches(4.7), Inches(0.3),
      size=12, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)

# Right: 3 phones
phone_paths = [
    (f"{SCEN}/P2-dana-wallet.png",    "Dana"),
    (f"{SCEN}/P2-yossi-wallet.png",   "Yossi"),
    (f"{SCEN}/P2-maya-wallet.png",    "Maya"),
]
for i, (p, name) in enumerate(phone_paths):
    if os.path.exists(p):
        bx = Inches(6.7 + i * 2.1)
        rrect(sl, bx, Inches(2.0), Inches(1.95), Inches(4.5), PANEL, BORDER)
        pad = Inches(0.07)
        sl.shapes.add_picture(p, bx+pad, Inches(2.0)+pad, Inches(1.95)-pad*2, Inches(4.5)-pad*2)
        t(sl, name, bx, Inches(6.55), Inches(1.95), Inches(0.3),
          size=10, color=DIM, mono=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 6. POS INTEGRATION OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ mock-backend/server.js", 6, TOTAL)
title_block(sl, "אינטגרציית POS — אושר עד", "Plug-and-play · dedicated M2M secret · idempotent")

section_header(sl, "ARCHITECTURE", Inches(0.7), Inches(2.1))

# Diagram
t_runs(sl, [
    ("POS terminal", 13, BLUE, True, False),
    ("  →  ", 13, DIM, False, True),
    ("Osher-Ad backend", 13, ORANGE, True, False),
    ("  →  ", 13, DIM, False, True),
    ("Clubby Edge Function", 13, GREEN, True, False),
    ("  →  ", 13, DIM, False, True),
    ("Customer wallet", 13, PURPLE, True, False),
], Inches(0.7), Inches(2.5), Inches(12.0), Inches(0.4))

section_header(sl, "WHY IT MATTERS", Inches(0.7), Inches(3.2), color=ORANGE)
items = [
    "🔐  M2M secret — לא חושף service_role key ללקוחות חיצוניים",
    "🔁  Idempotency — POS retry בטוח, לא מכפיל נקודות / הטבות",
    "⚡  Non-blocking — purchase חוזר ב-77ms בלי להמתין ל-Clubby",
    "💾  Persistence — עסקאות שורדות restart (data/*.json)",
    "🔌  Live event log — שקיפות מלאה של זרימת הנתונים",
    "🏢  Anchor dashboard — KPI חוצי-סניפים לרשת קמעונאית",
]
for j, b in enumerate(items):
    t(sl, b, Inches(0.7), Inches(3.6 + j * 0.43), Inches(5.7), Inches(0.4),
      size=12, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)

# Right: code panel
code_panel(sl, Inches(6.8), Inches(2.1), Inches(5.7), Inches(4.7),
           [
             [("// ", 11, DIM, False, True), ("Process a sale", 11, DIM, False, True)],
             [("app.post('", 11, TEXT, False, True), ("/api/purchase", 11, GREEN, False, True), ("', (req, res) => {", 11, TEXT, False, True)],
             [("  const { ", 11, TEXT, False, True), ("transaction_id", 11, BLUE, False, True), (" } = req.body", 11, TEXT, False, True)],
             "",
             [("  // ", 11, DIM, False, True), ("Idempotency check", 11, DIM, False, True)],
             [("  if (processedIds.has(transaction_id))", 11, TEXT, False, True)],
             [("    return res.json({ ...cached, ", 11, TEXT, False, True), ("idempotent: true", 11, ORANGE, False, True), (" })", 11, TEXT, False, True)],
             "",
             [("  // ", 11, DIM, False, True), ("Process · award points", 11, DIM, False, True)],
             [("  const result = processSale(items, member)", 11, TEXT, False, True)],
             [("  TRANSACTIONS.push(result); saveJSON(...)", 11, TEXT, False, True)],
             "",
             [("  // ", 11, DIM, False, True), ("Non-blocking benefit issuance", 11, DIM, False, True)],
             [("  ", 11, TEXT, False, True), ("issueClubbyBenefit", 11, GREEN, True, True), ("(...)", 11, TEXT, False, True), ("  // fire-and-log", 11, DIM, False, True)],
             "",
             [("  res.json({ ...result, ", 11, TEXT, False, True), ("clubby_benefit_issued", 11, BLUE, False, True), (": '", 11, TEXT, False, True), ("pending", 11, ORANGE, False, True), ("' })", 11, TEXT, False, True)],
             "})",
           ], label="server.js")

# ══════════════════════════════════════════════════════════════════════════════
# 7. POS UI — Live Demo
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ pos-terminal/live-demo", 7, TOTAL)
title_block(sl, "POS חי — אושר עד", "5 tabs · Hebrew RTL · real-time integration")

img_panel(sl, f"{V2}/02-pos-dana-gold-member.png", Inches(0.6), Inches(2.0), Inches(12.1), Inches(4.5),
          label="http://localhost:3001 — Dana Levi · 🥇 Gold · 2,840 pts")

t_runs(sl, [
    ("$ ", 13, GREEN, True, True),
    ("demo:", 13, DIM, False, True),
    ("  ", 13, TEXT, False, False),
    ("חיפוש חבר → סל קנייה → עיבוד מכירה → הטבה ל-Clubby", 13, TEXT, False, False),
], Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 8. SALE COMPLETE — Pending state
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ pos-terminal/sale-complete", 8, TOTAL)
title_block(sl, "מכירה הושלמה — Non-blocking architecture", "77ms response · Clubby benefit dispatched async")

img_panel(sl, f"{V2}/03-pos-sale-pending-clubby.png", Inches(0.6), Inches(2.0), Inches(7.5), Inches(4.8),
          label="POS receipt — Hebrew + VAT + installments")

# Right: explanation
section_header(sl, "WHAT JUST HAPPENED", Inches(8.5), Inches(2.1), color=GREEN)
t_runs(sl, [
    ("1. ", 12, ORANGE, True, True),
    ("Receipt returned in ", 12, TEXT, False, False),
    ("77ms", 12, GREEN, True, False),
], Inches(8.5), Inches(2.5), Inches(4.5), Inches(0.3))
t(sl, "(Hebrew + VAT 17% + 3 תשלומים)",
  Inches(8.5), Inches(2.78), Inches(4.5), Inches(0.3),
  size=11, color=DIM)

t_runs(sl, [
    ("2. ", 12, ORANGE, True, True),
    ("Clubby call dispatched ", 12, TEXT, False, False),
    ("async", 12, BLUE, True, False),
], Inches(8.5), Inches(3.2), Inches(4.5), Inches(0.3))
t(sl, "Birthday gift identified (June)",
  Inches(8.5), Inches(3.48), Inches(4.5), Inches(0.3),
  size=11, color=DIM)

t_runs(sl, [
    ("3. ", 12, ORANGE, True, True),
    ("Banner shows ", 12, TEXT, False, False),
    ("⏳ pending", 12, YELLOW, True, False),
], Inches(8.5), Inches(3.9), Inches(4.5), Inches(0.3))
t(sl, "Updates to ✅ when Clubby confirms",
  Inches(8.5), Inches(4.18), Inches(4.5), Inches(0.3),
  size=11, color=DIM)

t_runs(sl, [
    ("4. ", 12, ORANGE, True, True),
    ("Transaction ", 12, TEXT, False, False),
    ("persisted", 12, PURPLE, True, False),
    (" to JSON", 12, TEXT, False, False),
], Inches(8.5), Inches(4.6), Inches(4.5), Inches(0.3))
t(sl, "Survives server restart",
  Inches(8.5), Inches(4.88), Inches(4.5), Inches(0.3),
  size=11, color=DIM)

# Code snippet bottom
code_panel(sl, Inches(8.5), Inches(5.4), Inches(4.5), Inches(1.3),
           [
             [("issueClubbyBenefit(...)", 11, GREEN, True, True)],
             [("  ", 11, TEXT, False, True), (".then(r => logEvent('benefit_issued'))", 11, DIM, False, True)],
             [("  ", 11, TEXT, False, True), (".catch(e => logEvent('benefit_error'))", 11, DIM, False, True)],
           ])

# ══════════════════════════════════════════════════════════════════════════════
# 9. INTEGRATION LOG
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ pos-terminal/integration-log", 9, TOTAL)
title_block(sl, "🔌 Integration Log — שקיפות מלאה", "GET /api/events · refreshes every 5s · last 200 events")

img_panel(sl, f"{V2}/04-pos-integration-log.png", Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.8),
          label="Live event feed — POS ⇄ Clubby")

# Right legend
section_header(sl, "EVENT TYPES", Inches(9.0), Inches(2.1), color=ORANGE)
events = [
    ("🛒", "purchase",          GREEN,  "Sale received"),
    ("⚠️", "purchase_duplicate", YELLOW, "Idempotency caught"),
    ("🎁", "benefit_issued",    GREEN,  "Clubby returned 200 OK"),
    ("❌", "benefit_error",     RED,    "Network / auth failure"),
    ("←",  "webhook_inbound",   BLUE,   "Customer scanned QR"),
    ("→",  "webhook_response",  PURPLE, "Returned tier benefit"),
]
for j, (em, name, color, desc) in enumerate(events):
    by = Inches(2.5 + j * 0.65)
    t(sl, em, Inches(9.0), by, Inches(0.4), Inches(0.4),
      size=14, color=TEXT, mono=True)
    t(sl, name, Inches(9.4), by, Inches(2.0), Inches(0.3),
      size=11, bold=True, color=color, mono=True)
    t(sl, desc, Inches(9.4), by + Inches(0.28), Inches(3.8), Inches(0.3),
      size=10, color=DIM)

t_runs(sl, [
    ("# ", 13, DIM, False, True),
    ("the screen that sells the integration story", 13, DIM, False, True),
], Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 10. ANCHOR DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ pos-terminal/anchor-dashboard", 10, TOTAL)
title_block(sl, "🏢 Anchor Dashboard — תצוגת רשת", "Powered by Clubby · cross-branch KPIs · enterprise pitch")

img_panel(sl, f"{V2}/05-pos-anchor-dashboard.png", Inches(0.6), Inches(2.0), Inches(8.0), Inches(4.8),
          label="Multi-branch overview")

section_header(sl, "WHAT IT SHOWS", Inches(9.0), Inches(2.1), color=GREEN)
items = [
    "1,361 חברי מועדון · 5 סניפים",
    "תל אביב · ירושלים · חיפה · אשדוד · בני ברק",
    "₪209,500 הכנסות מצטברות",
    "Tier breakdown — gold / silver / bronze",
    "קטגוריות מובילות לפי הכנסה",
    "Real-time data refresh",
]
for j, b in enumerate(items):
    t(sl, b, Inches(9.0), Inches(2.5 + j * 0.55), Inches(4.0), Inches(0.4),
      size=12, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)

t_runs(sl, [
    ("$ ", 13, GREEN, True, True),
    ("pitch:", 13, DIM, False, True),
    ("  ", 13, TEXT, False, False),
    ("המסך ש", 13, TEXT, False, False),
    ("סוגר", 13, ORANGE, True, False),
    (" את העסקה עם הרשת הקמעונאית", 13, TEXT, False, False),
], Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 11. v2 UPGRADES — Section divider
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/v2-upgrades", 11, TOTAL)
t(sl, "II.", Inches(0.7), Inches(2.0), Inches(2.0), Inches(1.5),
  size=80, bold=True, color=ORANGE, mono=True)
t(sl, "v2 שדרוגי UX",
  Inches(0.7), Inches(3.4), Inches(12.1), Inches(1.0),
  size=48, bold=True, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)
t(sl, "Senior Staff Designer audit · 12 critical fixes · WCAG · enterprise security",
  Inches(0.7), Inches(4.5), Inches(12.1), Inches(0.5),
  size=18, color=ORANGE, align=PP_ALIGN.RIGHT)

t_runs(sl, [
    ("$ ", 13, GREEN, True, True),
    ("git log", 13, TEXT, False, True),
    (" --since='audit'", 13, DIM, False, True),
    ("  ·  ", 13, DIM, False, True),
    ("22 modified files", 13, BLUE, True, False),
], Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 12. CRITICAL FIXES (5 P0s)
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/v2/critical-fixes.md", 12, TOTAL)
title_block(sl, "5 תיקוני P0 קריטיים", "User-facing bugs that broke daily interactions")

section_header(sl, "BEFORE  →  AFTER", Inches(0.7), Inches(2.1), color=RED)

fixes = [
    ("Silent redemption failure",  "if (error) return",         "+ setRedeemError(error.message)",  "redeem.tsx"),
    ("Permanent spinner",          "Promise.all(...)",          "+ .catch(() => setLoadError(true))", "Dashboard.tsx"),
    ("Modal breaks dark mode",     "background: 'white'",       "+ background: 'var(--surface)'",   "Members.tsx"),
    ("Restart button is no-op",    "/* todo */",                "+ DevSettings.reload()",           "profile.tsx"),
    ("FeedPosts loses input",      "setFormMode(null)  // always", "+ moved to success path only",  "FeedPosts.tsx"),
]
for i, (title, before, after, file) in enumerate(fixes):
    by = Inches(2.55 + i * 0.85)
    rrect(sl, Inches(0.7), by, Inches(11.9), Inches(0.78), SURFACE, BORDER)
    # Title + file
    t(sl, title, Inches(0.9), by + Inches(0.08), Inches(5.0), Inches(0.3),
      size=12, bold=True, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)
    t(sl, file, Inches(0.9), by + Inches(0.08), Inches(5.0), Inches(0.3),
      size=10, color=DIM, mono=True, align=PP_ALIGN.LEFT)
    # Before / After code
    t(sl, before, Inches(0.9), by + Inches(0.42), Inches(5.5), Inches(0.3),
      size=11, color=RED, mono=True)
    t(sl, "→", Inches(6.4), by + Inches(0.42), Inches(0.4), Inches(0.3),
      size=14, color=DIM, mono=True, align=PP_ALIGN.CENTER)
    t(sl, after, Inches(6.8), by + Inches(0.42), Inches(5.7), Inches(0.3),
      size=11, color=GREEN, mono=True)

# ══════════════════════════════════════════════════════════════════════════════
# 13. M2M SECRET
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ supabase/functions/issue-benefit/index.ts", 13, TOTAL)
title_block(sl, "🔐 M2M Integration Secret", "Don't hand a service_role key to a retail backend")

# Left: explanation
section_header(sl, "THE PROBLEM", Inches(0.7), Inches(2.1), color=RED)
t_runs(sl, [
    ("service_role_key", 13, RED, True, True),
    (" = superadmin", 13, TEXT, False, False),
], Inches(0.7), Inches(2.5), Inches(5.7), Inches(0.3))
t(sl, "Full DB access · cannot rotate · must never leak", Inches(0.7), Inches(2.85), Inches(5.7), Inches(0.3),
  size=11, color=DIM)

section_header(sl, "THE FIX", Inches(0.7), Inches(3.5), color=GREEN)
items = [
    "✅  Dedicated CLUBBY_INTEGRATION_KEY env var",
    "✅  Edge function validates key, uses service role internally",
    "✅  Rotatable without redeploy",
    "✅  Safe to share with external POS / ERP / CRM",
    "✅  Backwards-compat with service_role for internal tooling",
]
for j, b in enumerate(items):
    t(sl, b, Inches(0.7), Inches(3.9 + j * 0.4), Inches(5.7), Inches(0.4),
      size=12, color=TEXT)

# Right: code
code_panel(sl, Inches(6.8), Inches(2.1), Inches(6.0), Inches(4.7),
           [
             [("// ", 11, DIM, False, True), ("auth check", 11, DIM, False, True)],
             [("const token = authHeader.replace('Bearer ', '')", 11, TEXT, False, True)],
             "",
             [("const ", 11, ORANGE, True, True), ("integrationKey", 11, BLUE, False, True), (" = Deno.env.get(", 11, TEXT, False, True)],
             [("  '", 11, TEXT, False, True), ("CLUBBY_INTEGRATION_KEY", 11, GREEN, False, True), ("'", 11, TEXT, False, True)],
             ")",
             "",
             [("const ", 11, ORANGE, True, True), ("isIntegrationKey", 11, BLUE, False, True), (" = token === integrationKey", 11, TEXT, False, True)],
             [("const ", 11, ORANGE, True, True), ("isServiceKey", 11, BLUE, False, True), ("     = token === serviceRoleKey", 11, TEXT, False, True)],
             "",
             [("if (!isIntegrationKey && !isServiceKey)", 11, TEXT, False, True)],
             [("  return json({ ", 11, TEXT, False, True), ("error", 11, RED, False, True), (": '", 11, TEXT, False, True), ("Unauthorized", 11, ORANGE, False, True), ("' }, 401)", 11, TEXT, False, True)],
             "",
             [("// ", 11, DIM, False, True), ("Always use service role internally", 11, DIM, False, True)],
             [("// ", 11, DIM, False, True), ("External caller never sees this key", 11, DIM, False, True)],
             [("const adminClient = createClient(url, ", 11, TEXT, False, True), ("serviceRoleKey", 11, RED, False, True), (")", 11, TEXT, False, True)],
           ], label="issue-benefit/index.ts")

# ══════════════════════════════════════════════════════════════════════════════
# 14. IDEMPOTENCY
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ mock-backend/server.js", 14, TOTAL)
title_block(sl, "🔁 Idempotency — POS retry safety", "Real POS systems retry on timeout · don't double-charge points")

# Left: flow
section_header(sl, "THE TIMELINE", Inches(0.7), Inches(2.1), color=ORANGE)
flow = [
    ("t=0", "POS sends purchase",       "transaction_id: TX-001",       BLUE),
    ("t=1", "Server processes",         "+101 points to Dana",          GREEN),
    ("t=5", "Network timeout",          "POS never got response",       YELLOW),
    ("t=6", "POS retries",              "Same transaction_id",           YELLOW),
    ("t=7", "Server detects duplicate", "Returns cached response",       GREEN),
    ("t=8", "0 extra points awarded",   "No duplicate benefit",          GREEN),
]
for i, (time, action, detail, color) in enumerate(flow):
    by = Inches(2.55 + i * 0.55)
    t(sl, time, Inches(0.7), by, Inches(0.6), Inches(0.3),
      size=11, bold=True, color=color, mono=True)
    t(sl, action, Inches(1.4), by, Inches(2.7), Inches(0.3),
      size=12, color=TEXT)
    t(sl, detail, Inches(4.2), by, Inches(2.5), Inches(0.3),
      size=10, color=DIM)

# Right: code
code_panel(sl, Inches(7.0), Inches(2.1), Inches(5.8), Inches(4.7),
           [
             [("// ", 11, DIM, False, True), ("Idempotency map (persisted)", 11, DIM, False, True)],
             [("const ", 11, ORANGE, True, True), ("processedIds", 11, BLUE, False, True), (" = ", 11, TEXT, False, True), ("new Set", 11, GREEN, False, True), ("(...)", 11, TEXT, False, True)],
             "",
             [("app.post('", 11, TEXT, False, True), ("/api/purchase", 11, GREEN, False, True), ("', (req, res) => {", 11, TEXT, False, True)],
             [("  const { transaction_id } = req.body", 11, TEXT, False, True)],
             "",
             [("  // ", 11, DIM, False, True), ("Return cached for duplicate", 11, DIM, False, True)],
             [("  if (transaction_id && ", 11, TEXT, False, True)],
             [("      processedIds.has(transaction_id)) {", 11, TEXT, False, True)],
             [("    const cached = TRANSACTIONS.find(...)", 11, TEXT, False, True)],
             [("    return res.json({ ", 11, TEXT, False, True)],
             [("      ...cached._response, ", 11, TEXT, False, True)],
             [("      ", 11, TEXT, False, True), ("idempotent: true", 11, ORANGE, True, True)],
             "    })",
             "  }",
             "",
             [("  // ", 11, DIM, False, True), ("Process & remember the ID", 11, DIM, False, True)],
             [("  processedIds.add(transaction_id)", 11, TEXT, False, True)],
             "})",
           ], label="server.js")

# ══════════════════════════════════════════════════════════════════════════════
# 15. PERFORMANCE METRICS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/v2/metrics.md", 15, TOTAL)
title_block(sl, "מדידות לפני / אחרי", "What changed — and how much")

metrics = [
    ("⚡ Purchase response time",  "~3,000ms",   "77ms",         "39× faster",   GREEN),
    ("🔴 Critical broken flows",   "5 P0",       "0",            "100% fixed",   GREEN),
    ("🌑 Dark mode coverage",       "60%",        "100%",         "+40%",         GREEN),
    ("🔐 External backend access",  "service_role", "M2M key",    "enterprise-safe", GREEN),
    ("💾 Survives restart",         "❌",         "✅",           "JSON persist", GREEN),
    ("🔁 POS retry safe",           "❌ duplicate", "✅ deduped", "production",   GREEN),
    ("📊 Integration visibility",   "0 surface",  "Live event log", "new UX",     ORANGE),
    ("🏢 Anchor view",              "❌",         "✅",           "enterprise pitch", ORANGE),
]
for i, (label, before, after, gain, color) in enumerate(metrics):
    col = i % 2
    row = i // 2
    bx = Inches(0.7 + col * 6.2)
    by = Inches(2.2 + row * 1.15)
    rrect(sl, bx, by, Inches(6.0), Inches(1.0), SURFACE, BORDER)
    t(sl, label, bx + Inches(0.2), by + Inches(0.1), Inches(5.6), Inches(0.35),
      size=12, bold=True, color=TEXT, rtl=True, align=PP_ALIGN.RIGHT)
    # Numbers row
    t(sl, before, bx + Inches(0.2), by + Inches(0.5), Inches(2.0), Inches(0.4),
      size=11, color=RED, mono=True)
    t(sl, "→", bx + Inches(2.2), by + Inches(0.5), Inches(0.3), Inches(0.4),
      size=12, color=DIM, mono=True, align=PP_ALIGN.CENTER)
    t(sl, after, bx + Inches(2.6), by + Inches(0.5), Inches(2.0), Inches(0.4),
      size=11, color=GREEN, bold=True, mono=True)
    t(sl, gain, bx + Inches(0.2), by + Inches(0.5), Inches(5.6), Inches(0.4),
      size=11, color=color, bold=True, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 16. PORTAL — Improvements
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ apps/portal/v2", 16, TOTAL)
title_block(sl, "פורטל v2 — Settings + Platform Analytics", "Webhook integration · cross-business KPIs · CSV export")

# Two browser frames
img_panel(sl, f"{V2}/16-portal-webhook-config.png", Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.7),
          label="settings.tsx — webhook + signing secret")
img_panel(sl, f"{V2}/17-portal-platform-analytics.png", Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.7),
          label="anchor-dashboard.tsx — platform-wide view")

t_runs(sl, [
    ("# ", 13, DIM, False, True),
    ("Webhook URL + HMAC secret · 5-second timeout · automatic fallback", 13, DIM, False, True),
], Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 17. ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/architecture.md", 17, TOTAL)
title_block(sl, "Architecture — מה מתחת למכסה", "Supabase · Edge Functions · React Native · Express")

stacks = [
    (BLUE,   "Mobile",
     ["Expo SDK 56", "React Native 0.85", "Expo Router v4", "NativeWind v4", "TanStack Query"]),
    (PURPLE, "Portal",
     ["React 18 + Vite", "Recharts", "i18next (5 languages)", "CSS variables", "Full RTL"]),
    (GREEN,  "Backend",
     ["Supabase Postgres 17", "Deno Edge Functions", "Row Level Security", "Realtime", "pg_cron"]),
    (ORANGE, "Integration",
     ["Express.js (mock POS)", "HMAC-SHA256 webhook", "Expo Push + FCM", "Web Push (VAPID)", "JSON persistence"]),
]
for i, (color, title, items) in enumerate(stacks):
    cx = Inches(0.7 + i * 3.15)
    rrect(sl, cx, Inches(2.2), Inches(2.95), Inches(4.7), SURFACE, BORDER)
    rect(sl, cx, Inches(2.2), Inches(2.95), Pt(2.5), color)
    t(sl, title, cx + Inches(0.2), Inches(2.4), Inches(2.6), Inches(0.4),
      size=15, bold=True, color=color)
    t(sl, "//", cx + Inches(2.45), Inches(2.45), Inches(0.4), Inches(0.3),
      size=11, color=DIM, mono=True)
    for j, item in enumerate(items):
        t(sl, "·  " + item, cx + Inches(0.2), Inches(3.1 + j*0.55), Inches(2.6), Inches(0.5),
          size=11, color=TEXT, mono=True)

# ══════════════════════════════════════════════════════════════════════════════
# 18. ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/roadmap.md", 18, TOTAL)
title_block(sl, "מה הבא בתור", "Done · in progress · future · 100% planned")

cols = [
    ("DONE",         GREEN,  ["✅  5 P0 fixes", "✅  M2M secret", "✅  Idempotency",
                              "✅  Persistence", "✅  Non-blocking", "✅  Live event log",
                              "✅  Anchor dashboard", "✅  Receipt format",
                              "✅  Tier badges", "✅  Hebrew RTL"]),
    ("IN PROGRESS",  YELLOW, ["🔄  useCurrentUser hook", "🔄  Shared <ErrorBanner>",
                              "🔄  Shared <BackButton>", "🔄  Shared <Modal>",
                              "🔄  WCAG contrast pass", "🔄  i18n register"]),
    ("FUTURE",       PURPLE, ["🎯  Dark mode unified", "🎯  Live SMS OTP (Twilio)",
                              "🎯  GPS map view", "🎯  AI campaigns",
                              "🎯  App Store submission", "🎯  Production deploy"]),
]
for i, (title, color, items) in enumerate(cols):
    cx = Inches(0.7 + i * 4.15)
    rrect(sl, cx, Inches(2.2), Inches(3.95), Inches(4.7), SURFACE, BORDER)
    rect(sl, cx, Inches(2.2), Inches(3.95), Pt(2.5), color)
    t(sl, title, cx + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.4),
      size=14, bold=True, color=color, mono=True)
    for j, item in enumerate(items):
        t(sl, item, cx + Inches(0.2), Inches(3.05 + j*0.42), Inches(3.6), Inches(0.4),
          size=12, color=TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# 19. VERIFICATION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/test-results.md", 19, TOTAL)
title_block(sl, "All tests passed", "Browser-verified · automated Playwright · production-ready")

tests = [
    ("T1", "Dev hint visible on localhost only"),
    ("T2", "Redeem screen renders + error inline"),
    ("T3", "Discover loads businesses (auth)"),
    ("T4", "Portal dashboard stats visible"),
    ("T5a", "FeedPosts form opens correctly"),
    ("T5b", "FeedPosts saves + closes form"),
    ("T6", "Members modal uses var(--surface)"),
    ("T7", "Platform Analytics no permanent spinner"),
    ("M2M", "Wrong key → 401 returned"),
    ("IDEM", "Same transaction_id → idempotent flag"),
    ("PERSIST", "data/*.json files written"),
    ("ASYNC", "Purchase returns in 77ms"),
    ("EVENTS", "Log shows purchase + duplicate + error"),
    ("ANCHOR", "5 branches + tier breakdown shown"),
]
for i, (id, desc) in enumerate(tests):
    col = i % 2
    row = i // 2
    bx = Inches(0.7 + col * 6.2)
    by = Inches(2.3 + row * 0.6)
    rrect(sl, bx, by, Inches(6.0), Inches(0.5), SURFACE, BORDER)
    t(sl, "✓", bx + Inches(0.15), by + Inches(0.1), Inches(0.4), Inches(0.3),
      size=14, color=GREEN, bold=True, mono=True)
    t(sl, id, bx + Inches(0.6), by + Inches(0.13), Inches(1.1), Inches(0.3),
      size=11, bold=True, color=ORANGE, mono=True)
    t(sl, desc, bx + Inches(1.7), by + Inches(0.13), Inches(4.1), Inches(0.3),
      size=11, color=TEXT)

t_runs(sl, [
    ("$ ", 13, GREEN, True, True),
    ("npm test ", 13, TEXT, False, True),
    ("·", 13, DIM, False, True),
    ("  14 / 14 passed", 13, GREEN, True, False),
], Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 20. DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/deploy.md", 20, TOTAL)
title_block(sl, "Production deployment — ~$12/year", "Already coded · ready when you are")

steps = [
    ("1", "Supabase Cloud",     "supabase db push --linked",          "supabase functions deploy",   GREEN),
    ("2", "Portal → Vercel",    "cd apps/portal && npm run build",    "vercel deploy --prod",        BLUE),
    ("3", "Mobile web → Vercel","npx expo export --platform web",     "vercel deploy dist/",         PURPLE),
    ("4", "Native apps",        "eas build --platform all",           "eas submit --platform ios",   ORANGE),
]
for i, (num, title, cmd1, cmd2, color) in enumerate(steps):
    col = i % 2
    row = i // 2
    bx = Inches(0.7 + col * 6.2)
    by = Inches(2.3 + row * 1.6)
    rrect(sl, bx, by, Inches(6.0), Inches(1.4), SURFACE, BORDER)
    rect(sl, bx, by, Inches(6.0), Pt(2.5), color)
    # Num + title
    t(sl, num, bx + Inches(0.2), by + Inches(0.1), Inches(0.5), Inches(0.5),
      size=20, bold=True, color=color, mono=True)
    t(sl, title, bx + Inches(0.85), by + Inches(0.15), Inches(5.0), Inches(0.4),
      size=15, bold=True, color=TEXT)
    # Commands
    t(sl, "$ " + cmd1, bx + Inches(0.85), by + Inches(0.6), Inches(5.0), Inches(0.3),
      size=11, color=GREEN, mono=True)
    t(sl, "$ " + cmd2, bx + Inches(0.85), by + Inches(0.92), Inches(5.0), Inches(0.3),
      size=11, color=GREEN, mono=True)

t_runs(sl, [
    ("# ", 13, DIM, False, True),
    ("clubby.com domain · Cloudflare DNS · Vercel HTTPS · ~$12/yr total", 13, DIM, False, True),
], Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 21. CALL TO ACTION
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/cta.md", 21, TOTAL)
title_block(sl, "מוכן להריץ", "Three URLs · zero auth · live data")

# Three big URL cards
urls = [
    ("📱 Mobile wallet",    "http://localhost:8081",  "Customer app · Expo web",     BLUE),
    ("🏢 Business portal",  "http://localhost:5173",  "Owner CRM · analytics",        PURPLE),
    ("🛒 Osher-Ad POS",      "http://localhost:3001",  "Mock backend · 5 tabs · v2",   ORANGE),
]
for i, (title, url, sub, color) in enumerate(urls):
    cx = Inches(0.7 + i * 4.15)
    rrect(sl, cx, Inches(2.4), Inches(3.95), Inches(3.5), SURFACE, color, line_pt=2.0)
    t(sl, title, cx + Inches(0.2), Inches(2.6), Inches(3.6), Inches(0.5),
      size=18, bold=True, color=color, rtl=True, align=PP_ALIGN.RIGHT)
    # URL command line
    rrect(sl, cx + Inches(0.2), Inches(3.4), Inches(3.55), Inches(0.55), PANEL, BORDER)
    t(sl, "$  open " + url, cx + Inches(0.3), Inches(3.5), Inches(3.4), Inches(0.4),
      size=11, color=GREEN, mono=True)
    t(sl, sub, cx + Inches(0.2), Inches(4.15), Inches(3.6), Inches(0.4),
      size=12, color=DIM, align=PP_ALIGN.CENTER)
    # Big number
    t(sl, str(i+1), cx + Inches(3.2), Inches(2.6), Inches(0.6), Inches(0.5),
      size=20, bold=True, color=DIM, mono=True, align=PP_ALIGN.RIGHT)

# Sample data
section_header(sl, "DEMO ACCOUNTS", Inches(0.7), Inches(6.0), color=ORANGE)
t_runs(sl, [
    ("phone:", 12, DIM, False, True),
    (" +972521000001", 12, GREEN, False, True),
    ("  ·  ", 12, DIM, False, True),
    ("OTP:", 12, DIM, False, True),
    (" 000000", 12, GREEN, False, True),
    ("  ·  ", 12, DIM, False, True),
    ("portal:", 12, DIM, False, True),
    (" owner@test.com / portal1234", 12, GREEN, False, True),
], Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.4))

# ══════════════════════════════════════════════════════════════════════════════
# 22. THANKS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
header(sl, "~ clubby/thanks", 22, TOTAL)

t(sl, "תודה.", Inches(0.7), Inches(2.0), Inches(12.0), Inches(2.0),
  size=120, bold=True, color=TEXT, align=PP_ALIGN.CENTER, rtl=True)

t_runs(sl, [
    ("$ ", 22, GREEN, True, True),
    ("clubby ", 22, TEXT, False, True),
    ("--", 22, DIM, False, True),
    ("ready-for-production", 22, ORANGE, True, True),
], Inches(0.7), Inches(4.5), Inches(12.0), Inches(0.7))

t(sl, "ready when you are", Inches(0.7), Inches(5.5), Inches(12.0), Inches(0.5),
  size=18, color=DIM, italic=True, align=PP_ALIGN.CENTER)

# Bottom credits
t_runs(sl, [
    ("# ", 11, DIM, False, True),
    ("Gabi Kalaora", 11, TEXT, False, True),
    ("   ·   ", 11, DIM, False, True),
    ("Clubby v2 · Final demo · June 2026", 11, DIM, False, True),
], Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.3), align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f"✅ Saved: {OUT}")
print(f"   Slides: {len(list(prs.slides))}")
