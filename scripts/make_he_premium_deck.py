"""
Clubby — Premium Hebrew Deck (Clubby_HE_Premium_2026.pptx)
Rich, exclusive dark design — gold accents, full RTL, real demo screenshots
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

BASE     = "/Users/I560647/Library/CloudStorage/OneDrive-SAPSE/Desktop/finance/Clubby"
D        = f"{BASE}/screenshots/demo"
S        = f"{BASE}/screenshots/scenario"
WORDMARK = f"{BASE}/logo/wordmark.png"
OUT      = f"{BASE}/slides/Clubby_HE_Premium_2026.pptx"
os.makedirs(f"{BASE}/slides", exist_ok=True)

# ── Premium Dark Palette ──────────────────────────────────────────────────────
BG      = RGBColor(0x08, 0x09, 0x0D)   # near-black
SURFACE = RGBColor(0x0F, 0x11, 0x17)   # card surface
CARD    = RGBColor(0x16, 0x19, 0x22)   # elevated card
BORDER  = RGBColor(0x22, 0x28, 0x38)   # subtle border
GREEN   = RGBColor(0x1A, 0x7A, 0x4A)   # brand green
GBRT    = RGBColor(0x2E, 0xCC, 0x71)   # bright green
GOLD    = RGBColor(0xD4, 0xAF, 0x37)   # premium gold
GOLD_LT = RGBColor(0xF5, 0xD8, 0x20)   # light gold highlight
SILVER  = RGBColor(0xA8, 0xB2, 0xC8)   # silver text
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF     = RGBColor(0xF8, 0xFA, 0xFC)
DIM     = RGBColor(0x4A, 0x52, 0x60)
LGRAY   = RGBColor(0x7A, 0x85, 0x99)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
TEAL    = RGBColor(0x06, 0xB6, 0xD4)
INDIGO  = RGBColor(0x63, 0x66, 0xF1)
ROSE    = RGBColor(0xF4, 0x3F, 0x5E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Helpers ──────────────────────────────────────────────────────────────────
def set_rtl(p_xml):
    pPr = p_xml.find(qn('a:pPr'))
    if pPr is None:
        pPr = etree.SubElement(p_xml, qn('a:pPr'))
        p_xml.insert(0, pPr)
    pPr.set('rtl', '1')

def fill_bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def t(slide, text, x, y, w, h, size=16, bold=False, color=WHITE,
      align=PP_ALIGN.RIGHT, italic=False, rtl=True):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    if rtl: set_rtl(p._p)

def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def rrect(slide, x, y, w, h, color, border_color=None):
    """Rounded rect shape=9"""
    s = slide.shapes.add_shape(9, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if border_color:
        s.line.color.rgb = border_color; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def img(slide, path, x, y, w, h=None):
    if not path or not os.path.exists(path): return False
    if h: slide.shapes.add_picture(path, x, y, w, h)
    else: slide.shapes.add_picture(path, x, y, w)
    return True

def gold_line(slide, x, y, w):
    """Thin gold separator line"""
    rect(slide, x, y, w, Pt(1.5), GOLD)

def wordmark(slide):
    if os.path.exists(WORDMARK):
        slide.shapes.add_picture(WORDMARK, Inches(0.28), Inches(0.18), Inches(2.2), Inches(0.58))
    else:
        t(slide, "clubby.", Inches(0.28), Inches(0.18), Inches(2.5), Inches(0.55),
          size=20, bold=True, color=GBRT, align=PP_ALIGN.LEFT, rtl=False)

def pgnum(slide, n, total):
    t(slide, f"{n} / {total}", Inches(0.25), H-Inches(0.38), Inches(1.2), Inches(0.3),
      size=9, color=DIM, align=PP_ALIGN.LEFT, rtl=False)

def gold_pill(slide, label, x, y, pw=Inches(1.6)):
    ph = Inches(0.3)
    p = slide.shapes.add_shape(9, x, y, pw, ph)
    p.fill.solid(); p.fill.fore_color.rgb = RGBColor(0x1A, 0x16, 0x08)
    p.line.color.rgb = GOLD; p.line.width = Pt(0.75)
    t(slide, label, x+Inches(0.08), y+Inches(0.04), pw-Inches(0.14), ph-Inches(0.06),
      size=9, bold=True, color=GOLD, align=PP_ALIGN.CENTER, rtl=False)

def phone_frame(slide, path, x, y, fw, fh):
    if not path or not os.path.exists(path): return
    frm = slide.shapes.add_shape(9, x, y, fw, fh)
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0x08, 0x0A, 0x10)
    frm.line.color.rgb = GOLD; frm.line.width = Pt(1.2)
    pad = Inches(0.09)
    slide.shapes.add_picture(path, x+pad, y+pad, fw-pad*2, fh-pad*2)

def browser_frame(slide, path, x, y, fw, fh):
    if not path or not os.path.exists(path): return
    chrome_h = Inches(0.24)
    frm = slide.shapes.add_shape(9, x, y, fw, fh)
    frm.fill.solid(); frm.fill.fore_color.rgb = RGBColor(0x0C, 0x0E, 0x16)
    frm.line.color.rgb = GOLD; frm.line.width = Pt(1.2)
    bar = slide.shapes.add_shape(1, x, y, fw, chrome_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(0x12, 0x15, 0x1F)
    bar.line.fill.background()
    for di, dc in enumerate([RGBColor(0xFC,0x5C,0x5C), RGBColor(0xFD,0xBC,0x40), RGBColor(0x27,0xC9,0x3F)]):
        d = slide.shapes.add_shape(9, x+Inches(0.1+di*0.18), y+Inches(0.075), Inches(0.09), Inches(0.09))
        d.fill.solid(); d.fill.fore_color.rgb = dc; d.line.fill.background()
    slide.shapes.add_picture(path, x, y+chrome_h, fw, fh-chrome_h)

def dark_card(slide, x, y, w, h, accent=None):
    c = slide.shapes.add_shape(9, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = CARD
    c.line.color.rgb = BORDER; c.line.width = Pt(0.75)
    if accent:
        rect(slide, x, y, w, Pt(2.5), accent)
    return c

def star_bg(slide, density=18):
    """Scatter tiny gold dots as a starfield"""
    import random; random.seed(42)
    for _ in range(density):
        sx = random.uniform(0.2, 13.0); sy = random.uniform(0.2, 7.2)
        sz = random.uniform(0.02, 0.06)
        d = slide.shapes.add_shape(9, Inches(sx), Inches(sy), Inches(sz), Inches(sz))
        alpha = random.randint(30, 90)
        d.fill.solid(); d.fill.fore_color.rgb = GOLD
        d.line.fill.background()

TOTAL = 28

# ══════════════════════════════════════════════════════════════════════════════
# 1. COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
star_bg(sl, 25)

# Side gold stripe
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, GOLD)

# Wordmark top-left
wordmark(sl)

# Central glow band
glow = rect(sl, 0, Inches(2.5), W, Inches(2.6), RGBColor(0x0A,0x14,0x0C))
glow.line.fill.background()

# Main headline — right aligned, very large
t(sl, "ארנק הנאמנות", Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.2),
  size=66, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
t(sl, "שמחבר הכל.", Inches(0.5), Inches(2.2), Inches(12.3), Inches(1.1),
  size=66, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(5.0), Inches(3.45), Inches(4.6))
t(sl, "לקוחות · עסקים · מערכות קיימות",
  Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.65),
  size=20, color=SILVER, align=PP_ALIGN.RIGHT)

# Bottom stat strip
rect(sl, 0, Inches(5.1), W, Inches(2.2), RGBColor(0x0C, 0x0F, 0x18))
gold_line(sl, 0, Inches(5.1), W)
stats = [
    ("📱", "ארנק לקוח", "iOS · Android · Web"),
    ("🏢", "פורטל עסקי", "CRM · אנליטיקה · פוסטים"),
    ("🔌", "אינטגרציה", "Webhook · REST API"),
    ("📊", "פלטפורמה", "תובנות רשת קמעונאית"),
]
for i, (em, title, sub) in enumerate(stats):
    x = Inches(0.5 + i * 3.2)
    rrect(sl, x, Inches(5.3), Inches(2.9), Inches(1.7), SURFACE, BORDER)
    t(sl, em + "  " + title, x+Inches(0.15), Inches(5.48), Inches(2.65), Inches(0.45),
      size=15, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    t(sl, sub, x+Inches(0.15), Inches(5.98), Inches(2.65), Inches(0.38),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 2. OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GOLD)
wordmark(sl); pgnum(sl, 2, TOTAL)

t(sl, "פלטפורמה אחת. שלושה מוצרים.", Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.75),
  size=34, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(5.5), Inches(1.75), Inches(5.1))
t(sl, "Clubby מחברת תוכניות נאמנות, מערכות קיימות וארנקי לקוחות בפלטפורמה אחת, מאובטחת ומעוצבת.",
  Inches(0.4), Inches(1.95), Inches(12.5), Inches(0.65),
  size=15, color=SILVER, align=PP_ALIGN.RIGHT)

cols = [
    (GBRT,   "📱", "ארנק לקוח",
     ["הרשמה עם טלפון — ללא סיסמה",
      "QR → הטבה בשתי שניות",
      "ארנק: חותמות, נקודות, דרגות",
      "פיד מבצעים בזמן אמת",
      "מצב לילה, עברית RTL"]),
    (GOLD,   "🏢", "פורטל עסקי",
     ["CRM מלא — חברים, הטבות, היסטוריה",
      "מבצעים, כרטיסיות, סטוריז",
      "התראות Push ממוקדות",
      "דשבורד אנליטיקה + ייצוא CSV",
      "פוסטים לפיד ארנק הלקוח"]),
    (TEAL,   "🔌", "אינטגרציה",
     ["Webhook עם חתימת HMAC-SHA256",
      "API לדחיפת הטבה מ-POS",
      "Fallback אוטומטי למבצעים",
      "מדריך אינטגרציה מובנה בפורטל",
      "מתחבר לכל מערכת קיימת"]),
]
for i, (accent, em, title, bullets) in enumerate(cols):
    cx = Inches(0.35 + i * 4.3)
    dark_card(sl, cx, Inches(2.9), Inches(4.1), Inches(4.2), accent)
    t(sl, em + "  " + title, cx+Inches(0.2), Inches(3.1), Inches(3.8), Inches(0.5),
      size=17, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    for j, b in enumerate(bullets):
        t(sl, "· " + b, cx+Inches(0.15), Inches(3.75+j*0.54), Inches(3.8), Inches(0.48),
          size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 3. SECTION — Customer App
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
star_bg(sl, 20)
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, GBRT)
# Number
t(sl, "01", Inches(0.35), Inches(0.5), Inches(2), Inches(0.65),
  size=13, bold=True, color=DIM, align=PP_ALIGN.LEFT, rtl=False)
# Glow behind text
glow = rect(sl, 0, Inches(2.2), W, Inches(2.2), RGBColor(0x07,0x12,0x0B))
glow.line.fill.background()
t(sl, "חוויית הלקוח", Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.8),
  size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
gold_line(sl, Inches(4.5), Inches(3.25), Inches(4.3))
t(sl, "מהסריקה הראשונה ועד לפרס — חלק, מהיר, ויוקרתי",
  Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.65),
  size=19, color=RGBColor(0xBB,0xF7,0xD0), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# HELPER: mobile slide (dark, phone right)
# ══════════════════════════════════════════════════════════════════════════════
def mobile_slide(n, title, subtitle, bullets, img_path, accent=GBRT, tag="לקוח"):
    sl = prs.slides.add_slide(blank)
    fill_bg(sl, BG)
    rect(sl, W-Inches(0.06), 0, Inches(0.06), H, accent)
    wordmark(sl); pgnum(sl, n, TOTAL)

    # Phone right side
    PH = Inches(6.4); PW = Inches(3.2)
    PX = W - PW - Inches(0.35); PY = (H - PH) / 2
    phone_frame(sl, img_path, PX, PY, PW, PH)

    # Text left
    TX = Inches(0.4); TW = Inches(9.2)
    gold_pill(sl, tag, Inches(0.4), Inches(0.72))
    t(sl, title, TX, Inches(1.18), TW, Inches(1.0),
      size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    gold_line(sl, TX, Inches(2.28), Inches(4.8))
    t(sl, subtitle, TX, Inches(2.45), TW-Inches(0.3), Inches(1.2),
      size=14, color=SILVER, align=PP_ALIGN.RIGHT)
    for j, b in enumerate(bullets):
        bx = TX; by = Inches(3.82 + j * 0.62)
        # Accent dot
        d = sl.shapes.add_shape(9, TW-Inches(0.15), by+Inches(0.1),
                                  Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = accent; d.line.fill.background()
        t(sl, b, bx, by, TW-Inches(0.35), Inches(0.55),
          size=13, color=SILVER, align=PP_ALIGN.RIGHT)
    return sl

def portal_slide(n, title, subtitle, bullets, img_path, accent=GOLD, tag="פורטל"):
    sl = prs.slides.add_slide(blank)
    fill_bg(sl, BG)
    rect(sl, W-Inches(0.06), 0, Inches(0.06), H, accent)
    wordmark(sl); pgnum(sl, n, TOTAL)

    FW = Inches(6.5); FH = Inches(5.9)
    browser_frame(sl, img_path, Inches(0.35), Inches(0.75), FW, FH)

    TX = Inches(7.2); TW = Inches(5.85)
    gold_pill(sl, tag, TX, Inches(0.72))
    t(sl, title, TX, Inches(1.18), TW, Inches(0.9),
      size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    gold_line(sl, TX+Inches(1.0), Inches(2.18), Inches(4.7))
    t(sl, subtitle, TX, Inches(2.35), TW, Inches(1.1),
      size=13, color=SILVER, align=PP_ALIGN.RIGHT)
    for j, b in enumerate(bullets):
        by = Inches(3.6 + j * 0.55)
        d = sl.shapes.add_shape(9, TX+TW-Inches(0.18), by+Inches(0.1),
                                  Inches(0.1), Inches(0.1))
        d.fill.solid(); d.fill.fore_color.rgb = accent; d.line.fill.background()
        t(sl, b, TX, by, TW-Inches(0.3), Inches(0.48),
          size=12, color=SILVER, align=PP_ALIGN.RIGHT)
    return sl

def section_slide(num_str, title, sub, accent):
    sl = prs.slides.add_slide(blank)
    fill_bg(sl, BG)
    star_bg(sl, 18)
    rect(sl, W-Inches(0.07), 0, Inches(0.07), H, accent)
    t(sl, num_str, Inches(0.35), Inches(0.5), Inches(2), Inches(0.65),
      size=13, bold=True, color=DIM, align=PP_ALIGN.LEFT, rtl=False)
    glow = rect(sl, 0, Inches(2.2), W, Inches(2.2), RGBColor(0x0A,0x0C,0x15))
    glow.line.fill.background()
    t(sl, title, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.8),
      size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    gold_line(sl, Inches(4.5), Inches(3.25), Inches(4.3))
    t(sl, sub, Inches(0.5), Inches(3.45), Inches(12.3), Inches(0.65),
      size=18, color=SILVER, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 4. SIGN-IN
# ══════════════════════════════════════════════════════════════════════════════
mobile_slide(4,
  "כניסה חכמה — ללא סיסמה",
  "הלקוח מזין מספר טלפון. אם הוא חבר מועדון — מגיע לארנק תוך שניות. טלפון לא מוכר מקבל הנחיה לסרוק QR של חנות.",
  ["📞 OTP לטלפון — ללא חשבון, ללא סיסמה",
   "🔴 שגיאה inline לטלפון לא מוכר",
   "⚡ סביבת פיתוח: OTP 000000 לבדיקות",
   "↩️  ניתוב: ארנק לחבר קיים / הרשמה לחדש"],
  f"{D}/11-mobile-signin-screen.png", GBRT, "אימות")

# ══════════════════════════════════════════════════════════════════════════════
# 5. WALLET
# ══════════════════════════════════════════════════════════════════════════════
mobile_slide(5,
  "ארנק — הכל במקום אחד",
  "יתרת קרדיט, מועדונים עם תג דרגה, פיד עדכונים חי מהעסקים, והטבות פעילות — ממשק יוקרתי ב-dark mode.",
  ["💰 יתרת קרדיט בשקלים עם היסטוריה",
   "🏪 MY CLUBS — כל מועדון עם תג ברונזה/כסף/זהב",
   "📢 UPDATES — פיד פוסטים בזמן אמת",
   "🎁 רשימת הטבות עם טאבים: הנחות / קרדיט / פריט חינם"],
  f"{S}/P2-dana-wallet.png", GBRT, "ארנק")

# ══════════════════════════════════════════════════════════════════════════════
# 6. DISCOVER + STORE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, TEAL)
wordmark(sl); pgnum(sl, 6, TOTAL)
PH = Inches(5.7); PW = Inches(2.75)
phone_frame(sl, f"{D}/10-mobile-dana-discover.png", Inches(0.3), Inches(0.88), PW, PH)
phone_frame(sl, f"{S}/P4-dana-store-member.png",    Inches(3.35), Inches(0.88), PW, PH)
TX = Inches(6.55); TW = Inches(6.55)
gold_pill(sl, "גילוי חנויות", TX, Inches(0.72))
t(sl, "גלה עסקים — לפי מיקום",
  TX, Inches(1.18), TW, Inches(0.8), size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(1.5), Inches(2.1), Inches(4.8))
t(sl, "עסקים ממוינים לפי מרחק GPS, סינון קטגוריה, פתוח עכשיו וחיפוש. לחיצה על עסק פותחת פרופיל מלא עם כיסוי, לוגו, שעות ומבצעים.",
  TX, Inches(2.25), TW, Inches(1.15), size=13, color=SILVER, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "📍 ממוין לפי מרחק כשמיקום אושר",
    "🌐 7 ריבועי קטגוריה צבעוניים",
    "🟢 סינון 'פתוח עכשיו' לפי שעות",
    "✓  תג הצטרפות — אין כניסה כפולה",
    "🏪 פרופיל: כיסוי, לוגו, מבצעים, שעות",
]):
    t(sl, b, TX, Inches(3.55+j*0.56), TW-Inches(0.3), Inches(0.5), size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 7. REDEEM
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 7, TOTAL)
PH = Inches(5.7); PW = Inches(2.75)
phone_frame(sl, f"{S}/P4-dana-wallet-benefit.png", Inches(0.3), Inches(0.88), PW, PH)
phone_frame(sl, f"{D}/14-mobile-redeem-screen.png", Inches(3.35), Inches(0.88), PW, PH)
TX = Inches(6.55); TW = Inches(6.55)
gold_pill(sl, "מימוש הטבה", TX, Inches(0.72))
t(sl, "סריקה. הטבה. מימוש.",
  TX, Inches(1.18), TW, Inches(0.8), size=28, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(1.5), Inches(2.1), Inches(4.8))
t(sl, "הלקוח סורק QR — ההטבה מגיעה לארנק תוך שתי שניות. מימוש בלחיצה אחת: מראה מסך לקופאי.",
  TX, Inches(2.25), TW, Inches(1.0), size=13, color=SILVER, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "📷 מצלמה נייטיבית + fallback ל-jsQR בדפדפן",
    "⚡ הטבה בארנק תוך 2 שניות",
    "🏷️ סוגי הטבה: הנחה / קרדיט / פריט חינם",
    "✅ מסך מימוש ברור — מראים לקופאי",
    "🔒 Webhook חתום HMAC — העסק שולט בהטבה",
]):
    t(sl, b, TX, Inches(3.45+j*0.58), TW-Inches(0.3), Inches(0.52), size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 8. LOYALTY MECHANICS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GOLD)
wordmark(sl); pgnum(sl, 8, TOTAL)
gold_pill(sl, "מנגנוני נאמנות", Inches(0.4), Inches(0.72), pw=Inches(2.0))
t(sl, "חותמות · נקודות · דרגות",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.75),
  size=32, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.0), Inches(2.05), Inches(5.6))
t(sl, "שלושה מנגנוני נאמנות במקביל — כל אחד מוסיף שכבת מוטיבציה לחזרה לחנות.",
  Inches(0.4), Inches(2.22), Inches(12.5), Inches(0.55),
  size=14, color=SILVER, align=PP_ALIGN.RIGHT)
mechs = [
    (AMBER, "🥇 כרטיסיות חותמות",
     ["בעל עסק מגדיר מספר חותמות ופרס",
      "הלקוח רואה התקדמות בארנק",
      "השלמה → הטבה אוטומטית",
      "הטבה נשלחת גם ל-Clubby"]),
    (SILVER, "🏅 דרגות נאמנות",
     ["הגדרת ספי ברונזה / כסף / זהב",
      "מבוסס על חותמות מצטברות",
      "תג דרגה על כרטיס המועדון",
      "התקדמות לדרגה הבאה בפרופיל"]),
    (TEAL,  "⭐ תוכנית נקודות",
     ["נקודות שנצברות בכל סריקת QR",
      "יתרה נפרדת לכל עסק",
      "קטלוג פרסים: הנחה / קרדיט / פריט",
      "לוח מובילים לחברים הטובים"]),
]
for i, (accent, title, bullets) in enumerate(mechs):
    cx = Inches(0.35 + i * 4.3)
    dark_card(sl, cx, Inches(3.05), Inches(4.15), Inches(4.0), accent)
    t(sl, title, cx+Inches(0.2), Inches(3.25), Inches(3.85), Inches(0.45),
      size=15, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    for j, b in enumerate(bullets):
        t(sl, "· " + b, cx+Inches(0.15), Inches(3.85+j*0.55), Inches(3.9), Inches(0.48),
          size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 9. MANUAL COUPON + FEED
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, INDIGO)
wordmark(sl); pgnum(sl, 9, TOTAL)
PH = Inches(5.7); PW = Inches(2.75)
phone_frame(sl, f"{D}/13-mobile-manual-coupon.png", Inches(0.3),  Inches(0.88), PW, PH)
phone_frame(sl, f"{S}/P6-dana-feed.png",             Inches(3.35), Inches(0.88), PW, PH)
TX = Inches(6.55); TW = Inches(6.55)
gold_pill(sl, "תכונות נוספות", TX, Inches(0.72), pw=Inches(1.9))
t(sl, "קופון ידני ופיד מבצעים",
  TX, Inches(1.18), TW, Inches(0.8), size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX+Inches(1.5), Inches(2.1), Inches(4.8))
t(sl, "קופון ידני",
  TX, Inches(2.3), TW, Inches(0.38), size=14, bold=True, color=INDIGO, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "שם חנות, כותרת, סוג (קרדיט/הנחה/פריט)",
    "ריבועי בחירה משנים צבע כפתור הפעיל",
    "שגיאות inline — ללא Alert",
    "מסומן כ'ידני' — לא מאומת על ידי החנות",
]):
    t(sl, "· " + b, TX, Inches(2.75+j*0.42), TW-Inches(0.2), Inches(0.38),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)
t(sl, "פיד עדכונים (בסגנון WhatsApp)",
  TX, Inches(4.65), TW, Inches(0.38), size=14, bold=True, color=INDIGO, align=PP_ALIGN.RIGHT)
for j, b in enumerate([
    "פוסטים מגיעים לכל חברי המועדון בזמן אמת",
    "4 סוגים: מבצע, הצעה, הכרזה, סטורי",
    "כפתור CTA עם טקסט מותאם אישית",
    "פג תוקף אוטומטי — RLS מסנן",
]):
    t(sl, "· " + b, TX, Inches(5.1+j*0.42), TW-Inches(0.2), Inches(0.38),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 10. SECTION — Business Portal
# ══════════════════════════════════════════════════════════════════════════════
section_slide("02", "הפורטל העסקי", "CRM · מבצעים · אנליטיקה · פוסטים · סטוריז", GOLD)

# ══════════════════════════════════════════════════════════════════════════════
# 11. PORTAL DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
portal_slide(11,
  "דשבורד — תמונה מלאה",
  "נתוני עסק בזמן אמת: חברים, מבצעים פעילים, הטבות שהונפקו, שיעור מימוש, גרפי 30 יום.",
  ["📊 4 כרטיסי סטטיסטיקה",
   "📈 גרף חברים חדשים — 30 יום",
   "📉 גרף מימושים יומיים",
   "⬇️  Export Analytics → CSV מלא",
   "🌐 בחירת עסק מרובה בסרגל",
  ], f"{D}/15-portal-dashboard.png", GOLD, "פורטל")

# ══════════════════════════════════════════════════════════════════════════════
# 12. MEMBERS CRM
# ══════════════════════════════════════════════════════════════════════════════
portal_slide(12,
  "CRM חברים מלא",
  "רשימת כל חברי המועדון עם שם, טלפון ותאריך הצטרפות. לחיצה פותחת מודל עם היסטוריית הטבות.",
  ["👥 שם, טלפון, תאריך לכל חבר",
   "🔍 View → מודל עם כל ההטבות",
   "📣 שלח הודעה — Push ממוקד",
   "🎯 קהלים: כולם / חדשים / פעילים / ישנים",
   "⬇️  ייצוא CSV לכלי CRM חיצוני",
  ], f"{D}/16-portal-members.png", GOLD, "פורטל")

# ══════════════════════════════════════════════════════════════════════════════
# 13. PROMOTIONS + STAMPS
# ══════════════════════════════════════════════════════════════════════════════
portal_slide(13,
  "מבצעים, חותמות ונקודות",
  "ניהול מלא של מכניקות הנאמנות: מבצעי קבלת פנים, כרטיסיות חותמות, ספי דרגה, תוכניות נקודות וקודי QR מרובי-סניפים.",
  ["🎁 CRUD מבצעים — עריכה, שכפול, מחיקה",
   "🥇 כרטיסיות חותמות — מספר ופרס",
   "🏅 דרגות — ספי ברונזה/כסף/זהב",
   "⭐ נקודות — קצב צבירה וקטלוג פרסים",
   "📱 QR רב-סניפי — QR נפרד לכל סניף",
  ], f"{S}/P1-owner-stamp-cards.png", GOLD, "פורטל")

# ══════════════════════════════════════════════════════════════════════════════
# 14. FEED + STORIES
# ══════════════════════════════════════════════════════════════════════════════
portal_slide(14,
  "פוסטים ופיד",
  "עסקים מפרסמים ישירות לארנק חברי המועדון — מבצעים, הכרזות וסטוריז — בניהול מלא מהפורטל.",
  ["📣 4 סוגי פוסטים עם תווית צבעונית",
   "🔘 כפתור CTA מותאם אישית",
   "📅 תפוגה אוטומטית — RLS מסנן",
   "⏸️  עצור / הפעל / מחק ללא רענון",
   "👁️  כל חבר רואה רק פוסטים של עסקיו",
  ], f"{D}/19-portal-feed-posts.png", GOLD, "פורטל")

# ══════════════════════════════════════════════════════════════════════════════
# 15. NOTIFY
# ══════════════════════════════════════════════════════════════════════════════
portal_slide(15,
  "התראות Push חכמות",
  "שלח הודעה לכל החברים או לקהל ממוקד — חדשים, פעילים, או חברים שלא חזרו ב-30 יום.",
  ["📣 'שלח ל-7 חברים' — לחיצה אחת",
   "🎯 קהלים: כולם / חדשים / פעילים / ישנים",
   "🔔 כל חבר מקבל הודעה ל-iOS / Android",
   "📱 Web push בנפרד (VAPID)",
   "🔄 Re-engagement אוטומטי — Cron יומי",
  ], f"{S}/P5-owner-notify-modal.png", GOLD, "פורטל")

# ══════════════════════════════════════════════════════════════════════════════
# 16. SECTION — Backend Integration
# ══════════════════════════════════════════════════════════════════════════════
section_slide("03", "אינטגרציה", "מתחבר לכל מערכת קיימת — בפחות משעה", AMBER)

# ══════════════════════════════════════════════════════════════════════════════
# 17. WEBHOOK SETTINGS
# ══════════════════════════════════════════════════════════════════════════════
portal_slide(17,
  "הגדרה ב-3 שדות",
  "מזינים כתובת Webhook וסוד חתימה. Clubby חותם כל בקשה עם HMAC-SHA256 כדי שתוכלו לאמת שמקורה ב-Clubby בלבד.",
  ["🔐 חתימת HMAC-SHA256 — X-Clubby-Signature",
   "⏱️  Timeout 5 שניות + fallback אוטומטי",
   "📋 Fallback: מבצעים פעילים → ללא הטבה",
   "📄 מדריך אינטגרציה מובנה בפורטל",
   "🧪 ngrok / localtunnel לפיתוח מקומי",
  ], f"{D}/17-portal-settings-webhook.png", AMBER, "אינטגרציה")

# ══════════════════════════════════════════════════════════════════════════════
# 18. OSHER-AD POS — LOOKUP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 18, TOTAL)
FW = Inches(8.2); FH = Inches(6.0)
browser_frame(sl, f"{D}/02-pos-member-lookup-dana.png", Inches(0.3), Inches(0.8), FW, FH)
gold_pill(sl, "הדגמה חיה — POS אושר עד", Inches(0.45), Inches(6.88), pw=Inches(2.5))
# Side labels
TX = Inches(8.9); TW = Inches(4.1)
gold_pill(sl, "חבר: Dana Levi", TX, Inches(1.1), pw=Inches(2.0))
for j, b in enumerate([
    "דרגה: 🥇 זהב",
    "נקודות: 2,840",
    "הוצאה שנתית: ₪5,680",
    "15 מוצרים בשמות עברי",
    "EAN-13 ברקודים",
    "בחירה מהירה: דנה / יוסי / מאיה",
]):
    t(sl, b, TX, Inches(1.65+j*0.56), TW, Inches(0.5), size=13, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 19. OSHER-AD POS — SALE COMPLETE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 19, TOTAL)
FW = Inches(8.2); FH = Inches(6.0)
browser_frame(sl, f"{D}/04-pos-sale-complete-clubby-benefit.png", Inches(0.3), Inches(0.8), FW, FH)
gold_pill(sl, "הדגמה חיה — מכירה הושלמה + הטבה ל-Clubby", Inches(0.45), Inches(6.88), pw=Inches(3.5))
TX = Inches(8.9); TW = Inches(4.1)
gold_pill(sl, "תוצאה", TX, Inches(1.1), pw=Inches(1.0))
for j, b in enumerate([
    "🎂 חודש יום הולדת → פריט חינם",
    "🥇 מכפיל זהב ×2 נקודות",
    "📋 קבלה בעברית עם מע\"מ",
    "💳 3 תשלומים (אשראי)",
    "✅ הטבה נשלחה ל-Clubby!",
]):
    t(sl, b, TX, Inches(1.7+j*0.62), TW, Inches(0.55), size=13, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 20. POS DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 20, TOTAL)
FW = Inches(8.2); FH = Inches(6.0)
browser_frame(sl, f"{D}/07-pos-dashboard.png", Inches(0.3), Inches(0.8), FW, FH)
TX = Inches(8.9); TW = Inches(4.1)
gold_pill(sl, "דשבורד POS", TX, Inches(0.72), pw=Inches(1.8))
t(sl, "דשבורד ניהולי\nPOS אושר עד",
  TX, Inches(1.18), TW, Inches(0.8), size=20, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX, Inches(2.1), TW)
for j, (val, label) in enumerate([
    ("2", "חברים פעילים היום"),
    ("389", "נקודות שהונפקו"),
    ("₪263", "הכנסות היום"),
    ("5", "חברי מועדון"),
]):
    by = Inches(2.4 + j * 0.8)
    dark_card(sl, TX, by, TW-Inches(0.1), Inches(0.68))
    t(sl, val, TX+Inches(0.15), by+Inches(0.05), Inches(1.2), Inches(0.55),
      size=20, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    t(sl, label, TX+Inches(1.45), by+Inches(0.18), TW-Inches(1.6), Inches(0.38),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)
t(sl, "טבלת עסקאות אחרונות + מובילי קניות",
  TX, Inches(5.75), TW, Inches(0.38), size=11, color=DIM, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 21. API ENDPOINTS
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, AMBER)
wordmark(sl); pgnum(sl, 21, TOTAL)
gold_pill(sl, "REST API", Inches(0.4), Inches(0.72), pw=Inches(1.2))
t(sl, "כל ה-Endpoints — במבט אחד",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.75),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.5), Inches(2.05), Inches(6.1))
endpoints = [
    ("GET",  "/api/members?phone=…",       "חיפוש חבר לפי טלפון — דרגה, נקודות, יתרה"),
    ("POST", "/api/purchase",              "עיבוד מכירה → צבירת נקודות → הטבה ל-Clubby"),
    ("GET",  "/api/members/:id/history",   "היסטוריית קניות — 10 האחרונות"),
    ("GET",  "/api/dashboard",             "הכנסות, מובילים, עסקאות היום"),
    ("GET",  "/api/products",              "קטלוג מוצרים (15 מוצרים עם שמות עברי)"),
    ("POST", "/clubby/webhook",            "↩ Clubby קורא לכאן בסריקת QR"),
    ("POST", "/functions/v1/issue-benefit","↩ דחיפת הטבה מ-POS לארנק הלקוח"),
]
for j, (method, path, desc) in enumerate(endpoints):
    by = Inches(2.6 + j * 0.63)
    mc = GBRT if method == "GET" else AMBER
    pill_w = Inches(0.7)
    pill_shape = sl.shapes.add_shape(9, Inches(0.4), by+Inches(0.02), pill_w, Inches(0.28))
    pill_shape.fill.solid(); pill_shape.fill.fore_color.rgb = RGBColor(0x0A,0x14,0x0C) if method == "GET" else RGBColor(0x14,0x0E,0x03)
    pill_shape.line.color.rgb = mc; pill_shape.line.width = Pt(0.75)
    t(sl, method, Inches(0.42), by+Inches(0.03), Inches(0.65), Inches(0.25),
      size=9, bold=True, color=mc, align=PP_ALIGN.CENTER, rtl=False)
    t(sl, path, Inches(1.22), by-Inches(0.02), Inches(4.5), Inches(0.32),
      size=12, bold=True, color=WHITE, align=PP_ALIGN.LEFT, rtl=False)
    t(sl, desc, Inches(6.0), by, Inches(6.9), Inches(0.3),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)
    rect(sl, Inches(0.4), by+Inches(0.52), Inches(12.5), Pt(0.5), BORDER)

# ══════════════════════════════════════════════════════════════════════════════
# 22. SECTION — Platform Analytics
# ══════════════════════════════════════════════════════════════════════════════
section_slide("04", "אנליטיקת פלטפורמה", "תובנות חוצות-עסקים ללקוח העוגן", TEAL)

# ══════════════════════════════════════════════════════════════════════════════
# 23. PLATFORM ANALYTICS
# ══════════════════════════════════════════════════════════════════════════════
portal_slide(23,
  "מבט פלטפורמאי",
  "דשבורד ייעודי לרשתות קמעונאיות — נתוני כל העסקים בתצוגה אחת: חברים, הטבות, שיעור מימוש וגרפי צמיחה.",
  ["🏢 4 עסקים · 12 חברים · 13 הטבות",
   "📈 גרפי משתמשים חדשים ב-30 יום",
   "🏆 טבלת עסקים מובילים לפי חברים",
   "📊 שיעור מימוש לכל עסק",
   "🔄 נתונים חיים — מתרענן בכל טעינה",
  ], f"{D}/18-portal-platform-analytics.png", TEAL, "עוגן")

# ══════════════════════════════════════════════════════════════════════════════
# 24. SCENARIO — 5 customers
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GBRT)
wordmark(sl); pgnum(sl, 24, TOTAL)
PH = Inches(5.5); PW = Inches(2.6)
phone_frame(sl, f"{S}/P2-dana-wallet.png",  Inches(0.25), Inches(0.97), PW, PH)
phone_frame(sl, f"{S}/P2-yossi-wallet.png", Inches(3.1),  Inches(0.97), PW, PH)
phone_frame(sl, f"{S}/P6-dana-feed.png",    Inches(5.95), Inches(0.97), PW, PH)
TX = Inches(9.2); TW = Inches(3.85)
gold_pill(sl, "תרחיש מלא", TX, Inches(0.72), pw=Inches(1.6))
t(sl, "5 לקוחות,\nבית קפה אחד",
  TX, Inches(1.18), TW, Inches(1.1), size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, TX, Inches(2.45), TW)
t(sl, "תרחיש קצה-לקצה שנבדק חי: בעל עסק מגדיר, 5 לקוחות נרשמים דרך QR, מקבלים הטבות, רואים פיד — והדשבורד מתעדכן בזמן אמת.",
  TX, Inches(2.62), TW, Inches(1.5), size=11, color=SILVER, align=PP_ALIGN.RIGHT)
for j, m in enumerate([
    "✅ QR → הטבה בפחות מ-2 שניות",
    "✅ תגי ברונזה/כסף/זהב",
    "✅ פוסט פיד → ארנק בזמן אמת",
    "✅ 'שלח ל-7 חברים'",
    "✅ אנליטיקה מתעדכנת",
]):
    t(sl, m, TX, Inches(4.3+j*0.52), TW, Inches(0.45), size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 25. BUILT vs ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GOLD)
wordmark(sl); pgnum(sl, 25, TOTAL)
gold_pill(sl, "מצב הפיתוח", Inches(0.4), Inches(0.72), pw=Inches(1.8))
t(sl, "מה בנוי · מה הבא",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.0), Inches(2.0), Inches(5.6))

built = [
    "✅ כניסה OTP לטלפון (לקוח + פורטל)",
    "✅ סריקת QR + הטבת קבלת פנים",
    "✅ ארנק: חותמות, נקודות, דרגות",
    "✅ Discover: GPS, קטגוריות, חיפוש",
    "✅ פורטל עסקי — CRM מלא + אנליטיקה",
    "✅ מבצעים, חותמות, דרגות, נקודות",
    "✅ סטוריז + פיד פוסטים",
    "✅ Push ממוקד (קהלים)",
    "✅ Webhook חתום HMAC",
    "✅ issue-benefit REST Edge Function",
    "✅ POS מדומה אושר עד — הדגמה חיה",
    "✅ אנליטיקת פלטפורמה (לקוח עוגן)",
    "✅ עברית RTL + i18n (5 שפות)",
    "✅ Dark mode + מערכת עיצוב Urbanist",
]
next_items = [
    "🔜 SMS אמיתי (Twilio/Vonage)",
    "🔜 תצוגת מפה ב-Discover",
    "🔜 סריקת ברקוד במצלמה",
    "🔜 דשבורד רשת קמעונאית",
    "🔜 קבלות דיגיטליות (דחיפה מ-POS)",
    "🔜 קמפיינים מבוססי AI",
    "🔜 שליחה לחנות אפליקציות",
    "🔜 הפעלה בייצור (Vercel + Supabase)",
]
# Built column
rect(sl, Inches(0.35), Inches(2.3), Inches(6.3), Inches(4.85), RGBColor(0x07, 0x12, 0x0B))
t(sl, "בנוי ועובד", Inches(0.55), Inches(2.42), Inches(5.9), Inches(0.38),
  size=13, bold=True, color=GBRT, align=PP_ALIGN.RIGHT)
for j, b in enumerate(built):
    t(sl, b, Inches(0.5), Inches(2.88+j*0.29), Inches(6.0), Inches(0.26),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)
# Roadmap column
rect(sl, Inches(6.9), Inches(2.3), Inches(6.1), Inches(4.85), RGBColor(0x12, 0x0E, 0x03))
t(sl, "הבא בתור", Inches(7.1), Inches(2.42), Inches(5.7), Inches(0.38),
  size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
for j, b in enumerate(next_items):
    t(sl, b, Inches(7.05), Inches(2.95+j*0.52), Inches(5.8), Inches(0.45),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 26. TECH STACK
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, SILVER)
wordmark(sl); pgnum(sl, 26, TOTAL)
gold_pill(sl, "טכנולוגיה", Inches(0.4), Inches(0.72), pw=Inches(1.2))
t(sl, "Stack טכנולוגי",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(5.5), Inches(2.0), Inches(5.1))
stacks = [
    (GBRT,  "📱 ארנק לקוח",
     "Expo SDK 56 + React Native 0.85\nExpo Router v4\nNativeWind v4 + Urbanist\nTanStack Query\nSupabase JS"),
    (GOLD,  "🏢 פורטל עסקי",
     "React 18 + Vite\nRecharts (גרפים)\nCSS variables\ni18next (5 שפות)\nRTL מלא"),
    (AMBER, "⚙️ Backend",
     "Supabase PostgreSQL 17\nDeno Edge Functions\nRow Level Security\nRealtime subscriptions\npg_cron"),
    (TEAL,  "🔌 אינטגרציה",
     "Node.js / Express\nHMAC-SHA256 Webhook\nExpo Push + FCM/APNs\nWeb Push (VAPID)\nSupabase Storage"),
]
for i, (accent, title, body) in enumerate(stacks):
    cx = Inches(0.35 + i * 3.25)
    dark_card(sl, cx, Inches(2.35), Inches(3.1), Inches(4.75), accent)
    t(sl, title, cx+Inches(0.18), Inches(2.52), Inches(2.85), Inches(0.42),
      size=14, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    t(sl, body, cx+Inches(0.15), Inches(3.05), Inches(2.88), Inches(3.8),
      size=12, color=SILVER, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 27. DEPLOYMENT
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
rect(sl, W-Inches(0.06), 0, Inches(0.06), H, GBRT)
wordmark(sl); pgnum(sl, 27, TOTAL)
gold_pill(sl, "פריסה לייצור", Inches(0.4), Inches(0.72), pw=Inches(1.8))
t(sl, "מוכן לייצור — ~$12 לשנה",
  Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.7),
  size=30, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.5), Inches(2.05), Inches(6.1))
t(sl, "כל מה שנדרש לפריסה כבר קוד. תוכנית הפריסה כתובה במלואה.",
  Inches(0.4), Inches(2.25), Inches(12.5), Inches(0.5),
  size=14, color=SILVER, align=PP_ALIGN.RIGHT)
steps = [
    (GBRT, "1", "Supabase Cloud",  "db push + Edge Functions\nTwilio credentials + pg_cron"),
    (GOLD, "2", "Vercel — פורטל", "portal.clubby.com\nnpm run build → vercel deploy"),
    (AMBER,"3", "Vercel — אפליקציה","app.clubby.com\nexpo export → vercel deploy"),
    (TEAL, "4", "חנויות האפליקציות","eas build → App Store\neas submit → Play Store"),
]
for i, (accent, num, title, body) in enumerate(steps):
    cx = Inches(0.35 + i * 3.25)
    dark_card(sl, cx, Inches(3.0), Inches(3.1), Inches(4.1), accent)
    circle = sl.shapes.add_shape(9, cx+Inches(2.5), Inches(3.12), Inches(0.48), Inches(0.48))
    circle.fill.solid(); circle.fill.fore_color.rgb = accent; circle.line.fill.background()
    t(sl, num, cx+Inches(2.55), Inches(3.15), Inches(0.38), Inches(0.38),
      size=15, bold=True, color=BG, align=PP_ALIGN.CENTER, rtl=False)
    t(sl, title, cx+Inches(0.15), Inches(3.15), Inches(2.35), Inches(0.45),
      size=13, bold=True, color=accent, align=PP_ALIGN.RIGHT)
    t(sl, body, cx+Inches(0.15), Inches(3.7), Inches(2.88), Inches(1.2),
      size=11, color=SILVER, align=PP_ALIGN.RIGHT)
t(sl, "דומיין: clubby.com (~$12/שנה Cloudflare) · HTTPS אוטומטי דרך Vercel",
  Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.32),
  size=11, color=DIM, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# 28. CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
fill_bg(sl, BG)
star_bg(sl, 30)
rect(sl, W-Inches(0.07), 0, Inches(0.07), H, GOLD)
wordmark(sl)

glow = rect(sl, 0, Inches(1.8), W, Inches(3.2), RGBColor(0x0A, 0x0F, 0x08))
glow.line.fill.background()

t(sl, "מוכנים לחבר", Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.1),
  size=52, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
t(sl, "את תוכנית הנאמנות שלכם.", Inches(0.5), Inches(2.35), Inches(12.3), Inches(1.05),
  size=52, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
gold_line(sl, W-Inches(6.0), Inches(3.5), Inches(5.6))
t(sl, "Full stack בנוי. אינטגרציה עובדת. הדגמה מוכנה.",
  Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.55),
  size=18, color=SILVER, align=PP_ALIGN.RIGHT)

for i, (label, val) in enumerate([
    ("ארנק לקוח", "http://localhost:8081"),
    ("פורטל עסקי", "http://localhost:5173"),
    ("POS אושר עד", "http://localhost:3001"),
]):
    bx = Inches(1.2 + i * 3.7)
    dark_card(sl, bx, Inches(4.8), Inches(3.5), Inches(1.9))
    t(sl, label, bx+Inches(0.15), Inches(4.9), Inches(3.25), Inches(0.38),
      size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
    t(sl, val, bx+Inches(0.15), Inches(5.35), Inches(3.25), Inches(0.35),
      size=10, color=DIM, align=PP_ALIGN.LEFT, rtl=False)

prs.save(OUT)
print(f"\n✅ Deck saved: {OUT}")
print(f"   Slides: {len(prs.slides)}")
